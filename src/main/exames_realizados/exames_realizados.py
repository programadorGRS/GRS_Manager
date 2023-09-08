import os
from datetime import datetime
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.util import Pt
from requests import Response
from sqlalchemy.exc import DatabaseError, SQLAlchemyError
from werkzeug.utils import secure_filename

from src import TIMEZONE_SAO_PAULO, UPLOAD_FOLDER, app
from src.email_connect import EmailConnect
from src.exporta_dados import ExportaDadosWS
from src.extensions import database as db
from src.main.empresa.empresa import Empresa
from src.main.empresa_principal.empresa_principal import EmpresaPrincipal
from src.main.exame.exame import Exame
from src.main.prestador.prestador import Prestador
from src.main.tipo_exame.tipo_exame import TipoExame
from src.main.unidade.unidade import Unidade
from src.utils import get_json_configs, zipar_arquivos


class ExamesRealizados(db.Model):
    __tablename__ = "ExamesRealizados"

    id_realizado = db.Column(db.Integer, primary_key=True)
    cod_empresa_principal = db.Column(
        db.Integer, db.ForeignKey("EmpresaPrincipal.cod"), nullable=False
    )
    seq_ficha = db.Column(db.Integer, nullable=False)
    cod_funcionario = db.Column(db.Integer, nullable=False)
    cpf = db.Column(db.String(30))
    nome_funcionario = db.Column(db.String(150))
    data_criacao = db.Column(db.Date, nullable=False)
    data_ficha = db.Column(db.Date, nullable=False)
    data_resultado = db.Column(db.Date, nullable=False)

    cod_tipo_exame = db.Column(
        db.Integer, db.ForeignKey("TipoExame.cod_tipo_exame"), nullable=False
    )
    id_prestador = db.Column(db.Integer, db.ForeignKey("Prestador.id_prestador"))
    id_empresa = db.Column(
        db.Integer, db.ForeignKey("Empresa.id_empresa"), nullable=False
    )
    id_unidade = db.Column(
        db.Integer, db.ForeignKey("Unidade.id_unidade"), nullable=False
    )
    id_exame = db.Column(db.Integer, db.ForeignKey("Exame.id_exame"), nullable=False)

    cod_setor = db.Column(db.String(255))
    nome_setor = db.Column(db.String(255))
    cod_cargo = db.Column(db.String(255))
    nome_cargo = db.Column(db.String(255))

    BASE_PPT = "src/main/exames_realizados/ppt/base_exames_realizados_v2.pptx"

    COLS_PLANILHA = [
        "cod_empresa_principal",
        "nome",
        "cod_empresa",
        "razao_social",
        "cod_unidade",
        "nome_unidade",
        "cod_setor",
        "nome_setor",
        "cod_cargo",
        "nome_cargo",
        "cod_funcionario",
        "cpf",
        "nome_funcionario",
        "seq_ficha",
        "cod_exame",
        "nome_exame",
        "cod_tipo_exame",
        "nome_tipo_exame",
        "data_criacao",
        "data_ficha",
        "data_resultado",
        "cod_prestador",
        "nome_prestador",
    ]

    @classmethod
    def buscar_exames_realizados(
        cls,
        cod_emp_princ: int | None = None,
        id_empresa: int | None = None,
        id_unidade: int | None = None,
        id_exame: int | None = None,
        data_inicio: datetime | None = None,
        data_fim: datetime | None = None,
    ):
        filtros = []

        if cod_emp_princ:
            filtros.append(cls.cod_empresa_principal == cod_emp_princ)
        if id_empresa:
            filtros.append(cls.id_empresa == id_empresa)
        if id_exame:
            filtros.append(cls.id_exame == id_exame)
        if data_inicio:
            filtros.append(cls.data_resultado >= data_inicio)
        if data_fim:
            filtros.append(cls.data_resultado <= data_fim)
        if id_unidade:
            filtros.append(cls.id_unidade == id_unidade)

        query = (
            db.session.query(  # type: ignore
                cls,
                EmpresaPrincipal.nome,
                Empresa.cod_empresa,
                Empresa.razao_social,
                Unidade.cod_unidade,
                Unidade.nome_unidade,
                Prestador.cod_prestador,
                Prestador.nome_prestador,
                TipoExame.cod_tipo_exame,
                TipoExame.nome_tipo_exame,
                Exame.cod_exame,
                Exame.nome_exame,
            )
            .outerjoin(
                EmpresaPrincipal, cls.cod_empresa_principal == EmpresaPrincipal.cod
            )
            .outerjoin(Empresa, cls.id_empresa == Empresa.id_empresa)
            .outerjoin(Unidade, cls.id_unidade == Unidade.id_unidade)
            .outerjoin(Prestador, cls.id_prestador == Prestador.id_prestador)
            .outerjoin(TipoExame, cls.cod_tipo_exame == TipoExame.cod_tipo_exame)
            .outerjoin(Exame, cls.id_exame == Exame.id_exame)
            .filter(*filtros)
        )

        return query

    @classmethod
    def inserir_exames_realizados(
        cls, id_empresa: int, data_inicio: datetime, data_fim: datetime
    ) -> dict[str, Any]:
        """
        Realiza request para o Exporta Dados Pedido Exame \
        e insere novos exames realizados na Tabela. Remove \
        exames ja realizados baseando no seq_ficha + id_exame.

        Returns:
            dict: {'status': str, 'qtd': int}
        """
        empresa: Empresa = Empresa.query.get(id_empresa)

        credenciais: dict = get_json_configs(
            empresa.empresa_principal.configs_exporta_dados
        )

        infos = {"status": "OK", "qtd": 0}

        par: dict = ExportaDadosWS.pedido_exame(
            empresa=empresa.cod_empresa,
            cod_exporta_dados=credenciais["PEDIDO_EXAMES_COD"],
            chave=credenciais["PEDIDO_EXAMES_KEY"],
            dataInicio=data_inicio.strftime("%d/%m/%Y"),
            dataFim=data_fim.strftime("%d/%m/%Y"),
        )

        response: dict = ExportaDadosWS.request_exporta_dados_ws(parametro=par)

        resp: Response = response["response"]
        if resp.status_code != 200:
            infos["status"] = "Erro no request"
            return infos

        erro_soc = response["erro_soc"]
        if erro_soc:
            msg_erro = response["msg_erro"]
            infos["status"] = msg_erro
            return infos

        df = ExportaDadosWS.xml_to_dataframe(xml_string=resp.text)

        if df.empty:
            infos["status"] = "df vazio"
            return infos

        df = cls.__tratar_df(df=df)

        # NOTE: manter apenas exames com data de resultado
        df = df[df["DATAEXAME"].notna()]

        if df.empty:
            infos["status"] = "df vazio"
            return infos

        df = cls.__get_infos_exame(df=df, cod_emp_princ=empresa.cod_empresa_principal)

        df = cls.__get_infos_unidade(df=df, id_empresa=empresa.id_empresa)

        df = cls.__get_infos_prestador(
            df=df, cod_emp_princ=empresa.cod_empresa_principal
        )

        df = cls.__remover_dados_atuais(df=df, id_empresa=empresa.id_empresa)

        if df.empty:
            infos["status"] = "df vazio"
            return infos

        try:
            qtd = cls.__inserir_exames_realizados(
                df=df,
                id_empresa=empresa.id_empresa,
                cod_emp_princ=empresa.cod_empresa_principal,
            )
            infos["qtd"] = qtd
        except (DatabaseError, SQLAlchemyError) as err:
            app.logger.error(err, exc_info=True)
            infos["status"] = str(err)

        return infos

    @staticmethod
    def __tratar_df(df: pd.DataFrame):
        df = df.copy()

        df = df.replace({"": None})

        COLS_INT = [
            "SEQUENCIAFICHA",
            "CODIGOFUNCIONARIO",
            "CODIGOTIPOEXAME",
            "CODIGOEMPRESA",
        ]

        for col in COLS_INT:
            df[col] = df[col].astype(int)

        COLS_DATE = [
            "DATACRIACAOPEDIDOEXAMES",
            "DATAFICHA",
            "DATAEXAME",
        ]

        for col in COLS_DATE:
            df[col] = pd.to_datetime(df[col], dayfirst=True, errors="coerce").dt.date

        # dropar linhas com data vazia
        df = df.dropna(axis=0, subset=COLS_DATE)

        df["CPFFUNCIONARIO"] = df["CPFFUNCIONARIO"].astype("string")

        return df

    @staticmethod
    def __get_infos_exame(df: pd.DataFrame, cod_emp_princ: int):
        df = df.copy()

        query = db.session.query(  # type: ignore
            Exame.id_exame,
            Exame.cod_exame,
        ).filter(Exame.cod_empresa_principal == cod_emp_princ)
        df_db = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        df = pd.merge(
            df,
            df_db,
            how="left",
            left_on="CODIGOINTERNOEXAME",
            right_on="cod_exame",
        )

        df = df.drop(columns="cod_exame")

        return df

    @staticmethod
    def __get_infos_unidade(df: pd.DataFrame, id_empresa: int):
        df = df.copy()

        query = db.session.query(Unidade.id_unidade, Unidade.cod_unidade).filter(  # type: ignore
            Unidade.id_empresa == id_empresa
        )
        df_db = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        df = pd.merge(
            df,
            df_db,
            how="left",
            left_on="CODIGOUNIDADE",
            right_on="cod_unidade",
        )

        df = df.drop(columns="cod_unidade")

        return df

    @staticmethod
    def __get_infos_prestador(df: pd.DataFrame, cod_emp_princ: int):
        df = df.copy()

        query = db.session.query(  # type: ignore
            Prestador.id_prestador, Prestador.cod_prestador
        ).filter(Prestador.cod_empresa_principal == cod_emp_princ)
        df_db = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        # NOTE: adicionar uma linha Nan no df db para nao travar o merge quando o prestador \
        # estiver vazio no df
        df_aux = pd.DataFrame({"id_prestador": [None], "cod_prestador": [None]})

        df_db = pd.concat([df_db, df_aux], axis=0, ignore_index=True)
        df_db = df_db.astype("Int32")

        df["CODIGOPRESTADOR"] = df["CODIGOPRESTADOR"].astype("Int32")

        df = pd.merge(
            df,
            df_db,
            how="left",
            left_on="CODIGOPRESTADOR",
            right_on="cod_prestador",
        )

        df = df.drop(columns="cod_prestador")

        return df

    @staticmethod
    def __remover_dados_atuais(df: pd.DataFrame, id_empresa: int):
        """Remover Exames Realizados que já existem na db"""
        df = df.copy()

        query = db.session.query(  # type: ignore
            ExamesRealizados.id_realizado,
            ExamesRealizados.seq_ficha,
            ExamesRealizados.id_exame,
        ).filter(ExamesRealizados.id_empresa == id_empresa)
        df_db = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        df = pd.merge(
            df,
            df_db,
            how="left",
            left_on=["SEQUENCIAFICHA", "id_exame"],
            right_on=["seq_ficha", "id_exame"],
        )

        # NOTE: manter apenas exames que não possuem id_realizado
        df = df[df["id_realizado"].isna()]

        df = df.drop(columns=["id_realizado", "seq_ficha"])

        return df

    @classmethod
    def __inserir_exames_realizados(
        cls, df: pd.DataFrame, id_empresa: int, cod_emp_princ: int
    ):
        COLS_RENAME = {
            "SEQUENCIAFICHA": "seq_ficha",
            "CODIGOFUNCIONARIO": "cod_funcionario",
            "CPFFUNCIONARIO": "cpf",
            "NOMEFUNCIONARIO": "nome_funcionario",
            "DATACRIACAOPEDIDOEXAMES": "data_criacao",
            "DATAFICHA": "data_ficha",
            "DATAEXAME": "data_resultado",
            "CODIGOTIPOEXAME": "cod_tipo_exame",
            "CODIGOSETOR": "cod_setor",
            "NOMESETOR": "nome_setor",
            "CODIGOCARGO": "cod_cargo",
            "NOMECARGO": "nome_cargo",
        }

        COLS_ID = ["id_prestador", "id_unidade", "id_exame"]

        COLS_FINAL = list(COLS_RENAME.values()) + COLS_ID

        df = df.copy()

        df = df.rename(columns=COLS_RENAME)

        df = df[COLS_FINAL]

        df["id_empresa"] = id_empresa
        df["cod_empresa_principal"] = cod_emp_princ

        # remover exames realizados sem unidades validas
        df = df.dropna(axis=0, subset=["id_unidade"])

        qtd = df.to_sql(
            name=cls.__tablename__,
            con=db.session.bind,  # type: ignore
            if_exists="append",
            index=False,
        )
        db.session.commit()  # type: ignore

        return qtd

    @classmethod
    def rotina_exames_realizados(
        cls,
        cod_empresa_principal: int,
        id_empresa: int,
        nome_empresa: str,
        data_inicio: datetime,
        data_fim: datetime,
        emails_destinatario: list[str],
        corpo_email: str,
        id_unidade: int | None = None,
        nome_unidade: str | None = None,
        filtro_exames: list[str] | None = None,
        ppt_trigger: int = 50,
        reply_to: list[str] | None = None
    ) -> dict[str, Any]:
        """
        Args:
            filtro_exames (list[str], optional): lista de ids dos Exames \
            para filtrar na query. Defaults to None.

            Gera ppt apenas se a qtd de linhas for maior que 50

        Returns:
            tuple[int, str]: (num de linhas da tabela, status email)
        """

        infos = {"status": "OK", "qtd": 0}

        query = ExamesRealizados.buscar_exames_realizados(
            cod_emp_princ=cod_empresa_principal,
            id_empresa=id_empresa,
            id_unidade=id_unidade,
            data_inicio=data_inicio,
            data_fim=data_fim,
        )

        if filtro_exames:
            query = query.filter(Exame.id_exame.in_(filtro_exames))

        df = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        if df.empty:
            infos["status"] = "df vazio"
            return infos

        infos["qtd"] = len(df)

        timestamp = int(datetime.now().timestamp())
        nome_entidade = secure_filename(nome_empresa).replace(".", "_")
        if nome_unidade:
            nome_entidade = secure_filename(nome_unidade).replace(".", "_")

        caminho_arqvs = os.path.join(
            UPLOAD_FOLDER, f"ExamesRealizados_{nome_entidade}_{timestamp}"
        )

        nome_excel = f"{caminho_arqvs}.xlsx"
        df_excel = df[ExamesRealizados.COLS_PLANILHA]
        df_excel.to_excel(nome_excel, index=False, freeze_panes=(1, 0))

        arquivos_zipar: list[str] = [nome_excel]

        if len(df) > ppt_trigger:
            nome_ppt = f"{caminho_arqvs}.pptx"

            ExamesRealizados.criar_ppt(
                df=df,
                nome_arquivo=nome_ppt,
                nome_empresa=nome_empresa,
                nome_unidade=nome_unidade,
            )

            arquivos_zipar.append(nome_ppt)

        pasta_zip = zipar_arquivos(
            caminhos_arquivos=arquivos_zipar,
            caminho_pasta_zip=f"{caminho_arqvs}.zip",
        )

        sbj = f"Exames Realizados Empresas - {nome_empresa}"
        if nome_unidade:
            sbj = f"Exames Realizados Unidades - {nome_unidade}"

        email_enviado = True
        erro_email = None
        try:
            EmailConnect.send_email(
                to_addr=emails_destinatario,
                message_subject=sbj,
                message_body=corpo_email,
                message_imgs=[EmailConnect.ASSINATURA_BOT],
                message_attachments=[pasta_zip],
                reply_to=reply_to,
            )
        except Exception as err:
            infos["status"] = f"erro ao enviar email {str(err)}"
            email_enviado = False
            erro_email = type(err).__name__

        log_email = EmailConnect(
            email_to=",".join(emails_destinatario),
            email_subject=sbj,
            attachments=pasta_zip,
            status=email_enviado,
            error=erro_email,
            df_len=len(df),
            email_date=datetime.now(tz=TIMEZONE_SAO_PAULO),
        )
        db.session.add(log_email)  # type: ignore
        db.session.commit()  # type: ignore

        return infos

    @classmethod
    def criar_ppt(
        cls,
        df: pd.DataFrame,
        nome_arquivo: str,
        nome_empresa: str,
        nome_unidade: str | None = None,
        data_origem: datetime = datetime.now(),
    ):
        # NOTE: ao editar um grafico de barras placeholder com varias series \
        # no Powerpoint, inserir series extras para evitar que o python crie \
        # series novas, assim as cores seguidas pelas series serao as do \
        # tema do PPT. Quando faltam series, as novas sao criadas com a mesma cor \
        # mas se tem series a mais, as que sobraram saem do PPT final

        df = df.copy()

        # instanciar apresentacao
        presentation = Presentation(cls.BASE_PPT)

        # SLIDE 1  - titulo--------------------------------------------------
        slide = presentation.slides[0]

        # periodo
        inicio = df["data_resultado"].min().strftime("%d/%m/%Y")
        fim = df["data_resultado"].max().strftime("%d/%m/%Y")
        slide.shapes[4].text_frame.paragraphs[0].runs[-1].text = f"{inicio} - {fim}"

        # nome empresa
        slide.shapes[5].text_frame.paragraphs[0].runs[-1].text = f"{nome_empresa}"

        # nome unidade
        if nome_unidade:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = f"{nome_unidade}"
        else:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = "Todas"

        # SLIDE 2 - indicadores gerais-------------------------------------------
        slide = presentation.slides[1]
        shapes = slide.shapes

        # total exames
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df))

        # total pedidos
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(
            len(df["seq_ficha"].drop_duplicates())
        )

        qtds = df[["seq_ficha", "nome_exame"]].groupby(by="seq_ficha").count()
        # media exames/pedido
        media = qtds.mean(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # max exames/pedido
        maximo = qtds.max(skipna=True).values[0]
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        # min exames/pedido
        minimo = qtds.min(skipna=True).values[0]
        shapes[7].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # media exames/unidade
        media = (
            df[["cod_unidade", "nome_exame"]]
            .groupby(by="cod_unidade")
            .count()
            .mean(skipna=True)
            .values[0]
        )
        shapes[8].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # media exames/setor
        media = (
            df[["cod_setor", "nome_exame"]]
            .groupby(by="cod_setor")
            .count()
            .mean(skipna=True)
            .values[0]
        )
        shapes[9].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # media exames/cargo
        media = (
            df[["cod_cargo", "nome_exame"]]
            .groupby(by="cod_cargo")
            .count()
            .mean(skipna=True)
            .values[0]
        )
        shapes[10].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # media exames/funcionario
        media = (
            df[["cod_funcionario", "nome_exame"]]
            .groupby(by="cod_funcionario")
            .count()
            .mean(skipna=True)
            .values[0]
        )
        shapes[11].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # SLIDE 3 - Exames por data---------------------------------------
        slide = presentation.slides[2]
        shapes = slide.shapes

        df["data_resultado"] = pd.to_datetime(df["data_resultado"])
        df["mesAno"] = df["data_resultado"].dt.strftime("%m/%Y")

        qtds = df[["mesAno", "nome_exame"]].groupby(by="mesAno").count()

        # media exames/data
        media = qtds.mean(skipna=True).values[0]
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # min exames/data
        minimo = qtds.min(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # max exames/data
        maximo = qtds.max(skipna=True).values[0]
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        df["mes"] = df["data_resultado"].dt.month
        df["ano"] = df["data_resultado"].dt.year

        dados = (
            df[["ano", "mes", "mesAno", "nome_tipo_exame", "nome_exame"]]
            .pivot_table(
                values="nome_exame",
                index=["ano", "mes", "mesAno"],
                columns="nome_tipo_exame",
                aggfunc="count",
            )
            .fillna(0)
        )
        dados.reset_index(level=["ano", "mes"], drop=True, inplace=True)
        # organizar colunas para manter a ordem no grafico
        dados.loc["Total"] = dados.sum(numeric_only=True)
        dados = dados.sort_values(by="Total", axis=1, ascending=False)
        dados.drop(labels=["Total"], axis=0, inplace=True)

        cls.editar_grafico(
            chart=slide.shapes[3].chart,
            categorias=list(dados.index),
            series={col: list(dados[col].values) for col in dados.columns},
        )

        # SLIDE 4 - exames por tipo ------------------------------------------------------
        slide = presentation.slides[3]
        shapes = slide.shapes

        qtds = (
            df[["nome_tipo_exame", "nome_exame"]].groupby(by="nome_tipo_exame").count()
        )
        # media
        media = qtds.mean(skipna=True).values[0]
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # min
        minimo = qtds.min(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # max
        maximo = qtds.max(skipna=True).values[0]
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        dados = df[["nome_tipo_exame", "nome_exame"]].groupby("nome_tipo_exame").count()
        dados["nome_exame"] = dados["nome_exame"] / dados["nome_exame"].sum()
        dados = dados.sort_values("nome_exame", axis=0, ascending=False)

        cls.editar_grafico(
            chart=slide.shapes[3].chart,
            categorias=list(dados.index),
            series={"qtd": list(dados["nome_exame"].values)},
        )

        # SLIDES DE TABELA 5-10 ------------------------------------------------------
        infos_slides_tabelas = [
            (4, "nome_unidade"),
            (5, "nome_setor"),
            (6, "nome_cargo"),
            (7, "nome_exame"),
            (8, "nome_prestador"),
            (9, "nome_funcionario"),
        ]

        for num_slide, nome_col in infos_slides_tabelas:
            slide = presentation.slides[num_slide]
            shapes = slide.shapes

            dados = df[[nome_col, "razao_social"]].groupby(nome_col).count().fillna(0)
            dados = dados.reset_index()
            dados = dados.sort_values("razao_social", ascending=False)

            cls.editar_tabela(shape=shapes[3].table, df=dados[:10])

        # SLIDE 11 referecias ------------------------------------------------------
        slide = presentation.slides[10]
        shapes = slide.shapes

        # relatorio gerado em
        slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = data_origem.strftime(
            "%d/%m/%Y"
        )

        presentation.save(nome_arquivo)
        return None

    @staticmethod
    def editar_grafico(chart, categorias: list, series: dict[list]):
        chart_data = ChartData()
        chart_data.categories = categorias
        if series:
            for nome_serie, vals in series.items():
                chart_data.add_series(nome_serie, vals)
        else:
            chart_data.add_series("Sem dados", [0])
        chart.replace_data(chart_data)
        return None

    @staticmethod
    def editar_tabela(shape: object, df: pd.DataFrame, font_size: int = 14):
        for coluna in range(len(df.columns)):
            # pular duas primeiras celulas de cada coluna
            dados_coluna = [None, None] + list(df[df.columns[coluna]].values)
            for linha in range(2, len(dados_coluna)):
                shape.cell(linha, coluna).text_frame.paragraphs[0].text = str(
                    dados_coluna[linha]
                )
                shape.cell(linha, coluna).text_frame.paragraphs[0].runs[
                    0
                ].font.size = Pt(font_size)
        return None
