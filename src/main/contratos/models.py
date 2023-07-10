import time
from datetime import datetime

import pandas as pd
from flask_sqlalchemy import BaseQuery
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.util import Pt
from requests.exceptions import HTTPError, RequestException
from werkzeug.utils import secure_filename

from manager import TIMEZONE_SAO_PAULO, UPLOAD_FOLDER, database
from manager.email_connect import EmailConnect
from manager.exporta_dados import ExportaDadosWS
from manager.models import Empresa, EmpresaPrincipal, Unidade
from manager.utils import get_json_configs, zipar_arquivos

# TODO: criar modelo do relatorio ppt
# TODO: criar funcao de criar relatorios xlsx e ppt
# TODO: finalizar tela para gerar relatorios via front end


class Contrato(database.Model):
    """
        Datas de Contratos entre GRS e Clientes

        Obs: id_unidade pode ser Null em casos onde o contrato é criado \
        no SOC sem ser associado a nenhuma Unidade. Para encontrar esses casos \
        no Exporta Dados, utilize o filtro de Unidade como 0 (zero)
    """
    __tablename__ = 'Contrato'
    id_contrato = database.Column(database.Integer, primary_key=True)
    cod_empresa_principal = database.Column(database.Integer, database.ForeignKey('EmpresaPrincipal.cod'), nullable=False)
    id_empresa = database.Column(database.Integer, database.ForeignKey('Empresa.id_empresa'), nullable=False)
    id_unidade = database.Column(database.Integer, database.ForeignKey('Unidade.id_unidade'))
    cod_produto = database.Column(database.Integer, nullable=False)
    nome_produto = database.Column(database.String(255))
    data_vencimento = database.Column(database.Date)
    data_realizado = database.Column(database.Date)
    situacao = database.Column(database.String(255))
    data_inclusao = database.Column(database.DateTime)

    BASE_PPT: str = 'manager/modules/contratos/ppt/base_contratos_v1.pptx'

    @classmethod
    def buscar_contratos(
        self,
        cod_empresa_principal: int | None = None,
        data_vencimento_inicio: datetime | None = None,
        data_vencimento_fim: datetime | None = None,
        data_realizado_inicio: datetime | None = None,
        data_realizado_fim: datetime | None = None,
        id_empresa: int | None = None,
        id_unidade: int | None = None,
        cod_produto: int | None = None,
        nome_produto: str | None = None,
        situacao: str | None = None
    ) -> BaseQuery:

        filtros = []

        if cod_empresa_principal:
            filtros.append(self.cod_empresa_principal == cod_empresa_principal)
        if data_vencimento_inicio:
            filtros.append(self.data_vencimento >= data_vencimento_inicio)
        if data_vencimento_fim:
            filtros.append(self.data_vencimento <= data_vencimento_fim)
        if data_realizado_inicio:
            filtros.append(self.data_realizado >= data_realizado_inicio)
        if data_realizado_fim:
            filtros.append(self.data_realizado <= data_realizado_fim)
        if id_empresa:
            filtros.append(self.id_empresa == id_empresa)
        if id_unidade != None:
            if id_unidade == 0:
                filtros.append(self.id_unidade == None)
            else:
                filtros.append(self.id_unidade == id_unidade)
        if cod_produto:
            filtros.append(self.cod_produto == cod_produto)
        if nome_produto:
            filtros.append(self.nome_produto.like(f'%{nome_produto}%'))
        if situacao:
            filtros.append(self.situacao == situacao)

        query: BaseQuery = (
            database.session.query(self)
            .filter(*filtros)
            .order_by(self.id_empresa, self.id_unidade, self.data_vencimento.desc())
        )
        return query

    @classmethod
    def inserir_contratos(self, id_empresa: int, id_unidade: int):
        """
        Realiza request para o Exporta dados Produtos com Data de Contrato [293]

        Trata DataFrame da Response, se houver e insere na database

        Returns:
            dict[str, any]: {
                        'cod_empresa_principal': int,
                        'nome_empresa_principal': str,
                        'id_empresa': int,
                        'nome_empresa': str,
                        'id_unidade': int,
                        'nome_unidade': str,
                        'status': str,
                        'qtd': int
                    }
        """
        empresa: Empresa = Empresa.query.get(id_empresa)
        empresa_principal: EmpresaPrincipal = EmpresaPrincipal.query.get(empresa.cod_empresa_principal)
        unidade: Unidade = Unidade.query.get(id_unidade)
        credenciais: dict[str, any] = get_json_configs(empresa_principal.configs_exporta_dados)

        infos: dict[str, any] = {
            'cod_empresa_principal': empresa_principal.cod,
            'nome_empresa_principal': empresa_principal.nome,
            'id_empresa': empresa.id_empresa,
            'nome_empresa': empresa.razao_social,
            'id_unidade': getattr(unidade, 'id_unidade', None),
            'nome_unidade': getattr(unidade, 'nome_unidade', None)
        }

        if id_unidade == 0:
            cod_unidade = 0 # se unidade for zero, buscar contratos sem unidade
        elif unidade:
            cod_unidade = unidade.cod_unidade

        parametro: dict[str, any] = ExportaDadosWS.produtos_data_contrato(
            cod_empresa_principal=empresa_principal.cod,
            cod_exporta_dados=credenciais['EXPORTADADOS_PROD_DATA_CONTRATO_COD'],
            chave=credenciais['EXPORTADADOS_PROD_DATA_CONTRATO_KEY'],
            empresaSelecao=empresa.cod_empresa,
            cod_unidade=cod_unidade
        )

        resp: dict[str, any] = {}
        for _ in range(3):
            try:
                resp = ExportaDadosWS.request_exporta_dados_ws(parametro=parametro, id_empresa=empresa.id_empresa)
                break
            except (HTTPError, RequestException) as error:
                time.sleep(3)
                infos['status'] = f'erro no request: {type(error).__name__}'
                infos['qtd'] = 0
                continue
        
        if not resp:
            return infos

        if resp['response'].status_code == 200:
            if not resp['erro_soc']:
                df = ExportaDadosWS.xml_to_dataframe(resp['response'].text)
                if not df.empty:
                    df.replace('', None, inplace=True)

                    for col in ['EMP', 'CODIGOPRODUTO']:
                        df.dropna(subset=col, axis=0, inplace=True)
                        df[col] = df[col].astype(int)
                    
                    df['DATAREALIZADO'] = pd.to_datetime(df['DATAREALIZADO'], dayfirst=True, errors='coerce').dt.date

                    # tratar coluna de vencimento
                    df[['day', 'month', 'year']] = df['DATAVENCIMENTO'].str.split('/', expand=True, n=3, regex=False).astype(int)
                    # zero significa ano atual (todos os anos)
                    df['year'].replace(0, datetime.now().year, inplace=True, regex=False)
                    df['data_vencimento'] = pd.to_datetime(arg=df[['day', 'month', 'year']]).dt.date
                    
                    df.rename(
                        columns={
                            'CODIGOPRODUTO': 'cod_produto',
                            'NOMEPRODUTO': 'nome_produto',
                            'DATAREALIZADO': 'data_realizado',
                            'SITUACAO': 'situacao'
                        },
                        inplace=True
                    )

                    df['cod_empresa_principal'] = empresa_principal.cod
                    df['id_empresa'] = empresa.id_empresa
                    df['id_unidade'] = getattr(unidade, 'id_unidade', None)
                    df['data_inclusao'] = datetime.now(tz=TIMEZONE_SAO_PAULO)

                    df = df[[
                        'cod_empresa_principal',
                        'id_empresa',
                        'id_unidade',
                        'cod_produto',
                        'nome_produto',
                        'data_vencimento',
                        'data_realizado',
                        'situacao',
                        'data_inclusao'
                    ]]

                    df.to_sql(
                        name=self.__tablename__,
                        con=database.session.bind,
                        index=False,
                        if_exists='append'
                    )
                    database.session.commit()
                    
                    infos['status'] = 'ok'
                    infos['qtd'] = len(df)
                else:
                    infos['status'] =  'vazio'
                    infos['qtd'] =  0
            else:
                infos['status'] = f"erro soc: {resp['msg_erro']}"
                infos['qtd'] = 0
        else:
            infos['status'] = 'erro no request'
            infos['qtd'] = 0

        return infos

    @classmethod
    def rotina_contratos(
        self,
        id_empresa: int,
        corpo_email: str,
        id_unidade: int | None = None,
        gerar_ppt: bool = True,
        testando: bool = True
    ) -> dict[str, any]:
        start = time.time()

        empresa: Empresa = Empresa.query.get(id_empresa)

        infos = {
            'cod_empresa_principal': empresa.empresa_principal.cod,
            'nome_empresa_principal': empresa.empresa_principal.nome,
            'id_empresa': empresa.id_empresa,
            'razao_social': empresa.razao_social,
            'emails': empresa.contratos_emails,
            'df_len': 0
        }

        if id_unidade:
            unidade: Unidade = Unidade.query.get(id_unidade)
            infos['id_unidade'] = unidade.id_unidade
            infos['nome_unidade'] = unidade.nome_unidade
            infos['emails'] = unidade.contratos_emails

        query: BaseQuery = (
            database.session.query(
                Contrato,
                EmpresaPrincipal.nome,
                Empresa.razao_social,
                Unidade.nome_unidade
            )
            .filter(Contrato.id_empresa == id_empresa)
            .join(EmpresaPrincipal, Contrato.cod_empresa_principal == EmpresaPrincipal.cod)
            .join(Empresa, Contrato.id_empresa == Empresa.id_empresa)
            .join(Unidade, Contrato.id_unidade == Unidade.id_unidade)
        )
        df: pd.DataFrame = pd.read_sql(sql=query.statement, con=database.session.bind)

        if not df.empty:
            infos['nome_arquivo'] = self.criar_relatorios(
                df=df,
                nome_empresa=empresa.razao_social,
                data_origem=df['data_inclusao'][0].date,
                nome_unidade=infos.get('nome_unidade', None),
                gerar_ppt=True
            )
            infos['df_len'] = len(df)

            if testando:
                enviar_para: list[str] = ['gabrielsantos@grsnucleo.com.br']
            else:
                enviar_para: list[str] = infos['emails'].split(';')
            
            assunto: str = f"Datas de Contratos - {empresa.razao_social}"

            try:
                EmailConnect.send_email(
                    to_addr=enviar_para,
                    reply_to=['gabrielsantos@grsnucleo.com.br', 'relacionamento@grsnucleo.com.br'],
                    message_subject=assunto,
                    message_body=corpo_email,
                    message_imgs=[EmailConnect.ASSINATURA_BOT],
                    message_attachments=[infos['nome_arquivo']]
                )
                status_email = True
                erro_email = None
                infos['status'] = 'OK'
            except Exception as erro:
                status_email = False
                erro_email = type(erro).__name__
                infos['status'] = 'Erro ao enviar email'
            
            # registrar envio
            log_email = EmailConnect(
                email_to=','.join(enviar_para),
                email_subject=assunto,
                attachments=infos['nome_arquivo'],
                status=status_email,
                error=erro_email,
                df_len=len(df),
                email_date=datetime.now(tz=TIMEZONE_SAO_PAULO)
            )
            database.session.add(log_email)
            database.session.commit()

        else:
            infos['nome_arquivo'] = None
            infos['status'] = 'Não enviado, vazio'

        infos['tempo_execucao'] = int(time.time() - start)

        return infos

    @classmethod
    def criar_relatorios(
        self,
        df: pd.DataFrame,
        nome_empresa: str,
        data_origem: datetime,
        nome_unidade: str = None,
        gerar_ppt: bool = True
    ) -> str:
        df['status'] = list(map(self.criar_status, df['data_vencimento']))

        # calcular dias ate o vencimento
        df['dias_vencer'] = (df['data_vencimento'] - datetime.now().date()).dt.days
        # substituir todos menores que 1 por NA
        df.loc[df['dias_vencer'] <= 0, 'dias_vencer'] = pd.NA

        # criar categorias a vencer
        df['a_vencer'] = pd.NA
        df.loc[df['dias_vencer'].between(1, 30, 'both'), 'a_vencer'] = 30
        df.loc[df['dias_vencer'].between(31, 60, 'both'), 'a_vencer'] = 60
        df.loc[df['dias_vencer'].between(61, 90, 'both'), 'a_vencer'] = 90
        df.loc[df['dias_vencer'].between(91, 365, 'both'), 'a_vencer'] = 365

        nome_arquivos = f"{UPLOAD_FOLDER}/Contratos_{secure_filename(nome_empresa).replace('.', '_')}_{int(datetime.now().timestamp())}"

        df = df[[
            'cod_empresa_principal',
            'nome',
            'id_empresa',
            'razao_social',
            'id_unidade',
            'nome_unidade',
            'cod_produto',
            'nome_produto',
            'data_vencimento',
            'data_realizado',
            'situacao',
            'status',
            'a_vencer',
        ]]
        df.to_excel(f'{nome_arquivos}.xlsx', index=False, freeze_panes=(1,0))

        arquivos_zipar = [f'{nome_arquivos}.xlsx']

        if gerar_ppt:
            self.criar_ppt(
                df=df,
                data_origem=data_origem,
                nome_arquivo=f'{nome_arquivos}.pptx',
                nome_empresa=nome_empresa,
                nome_unidade=nome_unidade
            )
            arquivos_zipar.append(f'{nome_arquivos}.pptx')

        pasta_zip: str = zipar_arquivos(
            caminhos_arquivos=arquivos_zipar,
            caminho_pasta_zip=f'{nome_arquivos}.zip'
        )

        return pasta_zip

    @classmethod
    def criar_ppt(
        self,
        df: pd.DataFrame,
        nome_arquivo: str,
        nome_empresa: str,
        nome_unidade: str = None,
        data_origem: datetime = datetime.now()
    ):
        # NOTE: ao editar um grafico de barras placeholder com varias series \
        # no Powerpoint, inserir series extras para evitar que o python crie \
        # series novas, assim as cores seguidas pelas series serao as do \
        # tema do PPT. Quando faltam series, as novas sao criadas com a mesma cor \
        # mas se tem series a mais, as que sobraram saem do PPT final

        df = df.copy()

        # instanciar apresentacao
        presentation = Presentation(self.BASE_PPT)


        # SLIDE 1  - titulo--------------------------------------------------
        slide = presentation.slides[0]

        # periodo
        inicio = df['data_vencimento'].min().strftime('%d/%m/%Y')
        fim = df['data_vencimento'].max().strftime('%d/%m/%Y')
        slide.shapes[4].text_frame.paragraphs[0].runs[-1].text = f'{inicio} - {fim}'

        # nome empresa
        slide.shapes[5].text_frame.paragraphs[0].runs[-1].text = f'{nome_empresa}'

        # nome unidade
        if nome_unidade:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = f'{nome_unidade}'
        else:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = 'Todas'


        # SLIDE 2 - indicadores gerais-------------------------------------------
        slide = presentation.slides[1]
        shapes = slide.shapes

        # total contratos
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df))

        # total a vencer
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df[df['status'] == 'A vencer']))

        # total vencidos
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df[df['status'] == 'Vencido']))

        # contratos por status
        dados = df[['status', 'cod_produto']].groupby(by='status').count()
        self.editar_grafico(
            chart=slide.shapes[6].chart,
            categorias=list(dados.index),
            series={'qtd': dados.values}
        )

        # contratos a vencer
        dados = df[['a_vencer', 'cod_produto']].groupby(by='a_vencer').count()
        self.editar_grafico(
            chart=slide.shapes[7].chart,
            categorias=list(dados.index),
            series={'qtd': dados.values}
        )


        # SLIDE 3 - Contratos por Unidade e Status---------------------------------------
        slide = presentation.slides[2]
        shapes = slide.shapes

        dados = (
            df[['nome_unidade', 'nome_produto', 'cod_produto', 'status']]
            .pivot_table(
                values='cod_produto',
                index=['nome_unidade', 'nome_produto'],
                columns='status',
                aggfunc='count',
            )
            .fillna(0)
        )
        dados['Total'] = dados.sum(axis=1, numeric_only=True)
        dados = dados.sort_values(by='Total', axis=0, ascending=False)
        dados.reset_index(inplace=True)

        # checar e organizar as colunas para a tabela do ppt
        for col in ['Ativo', 'A vencer', 'Vencido', 'Total']:
            if col in dados.columns:
                dados[col] = dados[col].astype(int)
            else:
                dados[col] = 0
        
        dados = dados[[
            'nome_unidade',
            'nome_produto',
            'Ativo',
            'A vencer',
            'Vencido',
            'Total'
        ]]

        self.editar_tabela(
            shape=shapes[3].table,
            df=dados[:10]
        )


        # SLIDE 4 - Contratos por Unidade e Vencimento -----------------------------
        slide = presentation.slides[3]
        shapes = slide.shapes

        dados = (
            df[['nome_unidade', 'nome_produto', 'cod_produto', 'a_vencer']]
            .pivot_table(
                values='cod_produto',
                index=['nome_unidade', 'nome_produto'],
                columns='a_vencer',
                aggfunc='count',
            )
            .fillna(0)
        )
        dados['Total'] = dados.sum(axis=1, numeric_only=True)
        dados = dados.sort_values(by='Total', axis=0, ascending=False)
        dados.reset_index(inplace=True)

        # checar e organizar as colunas para a tabela do ppt
        for col in [30, 60, 90, 365, 'Total']:
            if col in dados.columns:
                dados[col] = dados[col].astype(int)
            else:
                dados[col] = 0
        
        dados = dados[[
            'nome_unidade',
            'nome_produto',
            30,
            60,
            90,
            365,
            'Total'
        ]]

        self.editar_tabela(
            shape=shapes[3].table,
            df=dados[:10]
        )

        # SLIDE 5 referecias ------------------------------------------------------
        slide = presentation.slides[4]
        shapes = slide.shapes

        # relatorio gerado em
        slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = data_origem.strftime('%d/%m/%Y')

        presentation.save(nome_arquivo)
        return None

    @staticmethod
    def editar_grafico(
        chart,
        categorias: list,
        series: dict[list]
    ):
        chart_data = ChartData()
        chart_data.categories = categorias
        if series:
            for nome_serie, vals in series.items():
                chart_data.add_series(nome_serie, vals)
        else:
            chart_data.add_series('Sem dados', [0])
        chart.replace_data(chart_data)
        return None

    @staticmethod
    def editar_tabela(
        shape: object,
        df: pd.DataFrame,
        font_size: int=14
    ):
        for coluna in range(len(df.columns)):
            # pular duas primeiras celulas de cada coluna
            dados_coluna = [None, None] + list(df[df.columns[coluna]].values)
            for linha in range(2, len(dados_coluna)):
                shape.cell(linha, coluna).text_frame.paragraphs[0].text = str(dados_coluna[linha])
                shape.cell(linha, coluna).text_frame.paragraphs[0].runs[0].font.size = Pt(font_size)
        return None

    @staticmethod
    def criar_status(data_vencimento: datetime) -> str:
        today = datetime.now().date()
        if data_vencimento:
            if data_vencimento <= today:
                return 'Vencido'
            elif (data_vencimento - today).days <= 365:
                return 'A vencer'
            else:
                return 'Ativo'
        else:
            return 'Ativo'


