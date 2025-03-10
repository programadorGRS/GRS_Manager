import os
import time
from datetime import datetime

import pandas as pd
from flask_sqlalchemy import BaseQuery
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.util import Pt
from werkzeug.utils import secure_filename

from src import UPLOAD_FOLDER
from src.email_connect import EmailConnect
from src.exporta_dados import ExportaDadosWS
from src.extensions import database as db
from src.main.empresa.empresa import Empresa
from src.main.empresa_principal.empresa_principal import EmpresaPrincipal
from src.main.exame.exame import Exame
from src.main.funcionario.funcionario import Funcionario
from src.main.unidade.unidade import Unidade
from src.utils import get_json_configs, zipar_arquivos

from .ped_proc import PedidoProcessamento


class ConvExames(db.Model):
    __tablename__ = "ConvExames"
    id_conv = db.Column(db.Integer, primary_key=True)
    id_proc = db.Column(
        db.Integer, db.ForeignKey("PedidoProcessamento.id_proc"), nullable=False
    )
    cod_empresa_principal = db.Column(
        db.Integer, db.ForeignKey("EmpresaPrincipal.cod"), nullable=False
    )
    id_empresa = db.Column(
        db.Integer, db.ForeignKey("Empresa.id_empresa"), nullable=False
    )
    id_unidade = db.Column(
        db.Integer, db.ForeignKey("Unidade.id_unidade"), nullable=False
    )
    id_funcionario = db.Column(
        db.Integer, db.ForeignKey("Funcionario.id_funcionario"), nullable=False
    )
    id_exame = db.Column(db.Integer, db.ForeignKey("Exame.id_exame"), nullable=False)
    periodicidade = db.Column(db.Integer)
    data_adm = db.Column(db.Date)
    ult_pedido = db.Column(db.Date)
    data_res = db.Column(db.Date)
    refazer = db.Column(db.Date)

    BASE_PPT = "src/main/conv_exames/ppt/base_conv_exames_v2.pptx"

    COLUNAS_PLANILHA = [
        "cod_empresa",
        "razao_social",
        "unidade",
        "setor",
        "cargo",
        "cod_funcionario",
        "matricula",
        "cpf",
        "nome_funcionario",
        "cod_exame",
        "nome_exame",
        "periodicidade",
        "data_adm",
        "ult_pedido",
        "data_res",
        "refazer",
    ]

    STATUS_EXAMES = {
        1: "Exame Em dia",
        2: "Exame A vencer",
        3: "Exame Vencido",
        4: "Resultado do Exame Pendente",
        5: "Exame nunca foi realizado",
    }

    DIAS_VENCER_EXAMES = {0: "Vazio", 30: 30, 60: 60, 90: 90, 365: 365}

    PPT_TRIGGER = 50

    @classmethod
    def inserir_conv_exames(cls, id_proc: int) -> dict:
        """
        Consulta pedido de processamento e insere na db \
        se houver retorno. Registra o retorno na tabela PedidoProcessamento.

        Args:
            id_proc (int): id do PedidoProcessamento

        Returns:
            dict: {
                'cod_empresa_principal': int,
                'nome_empresa_principal': str,
                'id_empresa': int,
                'nome_empresa': str,
                'cod_solicitacao': int,
                'status': str,
                'qtd': int
            }
        """
        PED_PROC: PedidoProcessamento = PedidoProcessamento.query.get(id_proc)
        EMPRESA_PRINCIPAL: EmpresaPrincipal = EmpresaPrincipal.query.get(
            PED_PROC.cod_empresa_principal
        )
        EMPRESA: Empresa = Empresa.query.get(PED_PROC.id_empresa)
        CREEDENCIAIS: dict = get_json_configs(EMPRESA_PRINCIPAL.configs_exporta_dados)

        PARAMETRO: dict = ExportaDadosWS.consulta_conv_exames_assync(
            cod_empresa_principal=PED_PROC.cod_empresa_principal,
            cod_exporta_dados=CREEDENCIAIS["CONV_EXAMES_ASSYNC_COD"],
            chave=CREEDENCIAIS["CONV_EXAMES_ASSYNC_KEY"],
            cod_empresa_trab=PED_PROC.cod_empresa,
            cod_sol=PED_PROC.cod_solicitacao,
        )

        infos = {
            "cod_empresa_principal": EMPRESA_PRINCIPAL.cod,
            "nome_empresa_principal": EMPRESA_PRINCIPAL.nome,
            "id_empresa": EMPRESA.id_empresa,
            "nome_empresa": EMPRESA.razao_social,
            "cod_solicitacao": PED_PROC.cod_solicitacao,
            "status": None,
            "qtd": 0,
        }

        resp: dict = ExportaDadosWS.request_exporta_dados_ws(parametro=PARAMETRO)

        if resp["response"].status_code != 200:
            infos["status"] = "Erro no request"
            PED_PROC.obs = "Erro no request"
            db.session.commit()  # type: ignore
            return infos

        if resp["erro_soc"]:
            infos["status"] = f"erro soc: {resp['msg_erro']}"
            PED_PROC.obs = f'Erro soc: {resp["msg_erro"]}'
            db.session.commit()  # type: ignore
            return infos

        df = ExportaDadosWS.xml_to_dataframe(resp["response"].text)
        if df.empty:
            infos["status"] = "Vazio"
            PED_PROC.obs = "Vazio"
            PED_PROC.resultado_importado = True
            db.session.commit()  # type: ignore
            return infos

        df = ConvExames.tratar_df_conv_exames(
            df=df,
            cod_empresa_principal=PED_PROC.cod_empresa_principal,
            id_proc=PED_PROC.id_proc,
            id_empresa=PED_PROC.id_empresa,
        )

        qtd = df.to_sql(
            name=ConvExames.__tablename__,
            con=db.session.bind,  # type: ignore
            index=False,
            if_exists="append",
        )

        infos["status"] = "Ok"
        infos["qtd"] = qtd
        PED_PROC.resultado_importado = True
        PED_PROC.obs = "Consulta inserida"
        db.session.commit()  # type: ignore
        return infos

    @classmethod
    def tratar_df_conv_exames(
        cls, df: pd.DataFrame, cod_empresa_principal, id_proc: int, id_empresa: int
    ) -> pd.DataFrame:
        """
        Recebe dataframe de convocacao de exames

        Trata o dataframe

        Retorna dataframe
        """
        # validar possiveis colunas faltantes
        for col in ["ULTIMOPEDIDO", "DATARESULTADO", "REFAZER"]:
            if col not in df.columns:
                df[col] = None

        df.replace("", None, inplace=True)

        df.dropna(
            axis=0,
            subset=["CODIGOEMPRESA", "UNIDADE", "CODIGOFUNCIONARIO", "CODIGOEXAME"],
            inplace=True,
        )

        for col in ["CODIGOEMPRESA", "CODIGOFUNCIONARIO", "PERIODICIDADE"]:
            df[col] = df[col].fillna(0).astype(int)

        # converter cols de data
        for col in ["DATAADMISSAO", "ULTIMOPEDIDO", "DATARESULTADO", "REFAZER"]:
            df[col] = pd.to_datetime(
                df[col], dayfirst=True, errors="coerce"  # substitui erros por NaT
            ).dt.date

        df["id_proc"] = id_proc
        df["cod_empresa_principal"] = cod_empresa_principal

        # buscar ids
        query = db.session.query(  # type: ignore
            Funcionario.id_funcionario,
            Funcionario.cod_funcionario,
            Funcionario.id_empresa,
            Funcionario.id_unidade,
        ).filter(Funcionario.id_empresa == id_empresa)
        df_database = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        df = df.merge(
            df_database,
            how="left",
            left_on="CODIGOFUNCIONARIO",
            right_on="cod_funcionario",
        )

        query = db.session.query(  # type: ignore
            Exame.id_exame, Exame.cod_exame
        ).filter(Exame.cod_empresa_principal == cod_empresa_principal)
        df_database = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        df = df.merge(
            df_database, how="left", left_on="CODIGOEXAME", right_on="cod_exame"
        )

        df.rename(
            columns={
                "PERIODICIDADE": "periodicidade",
                "DATAADMISSAO": "data_adm",
                "ULTIMOPEDIDO": "ult_pedido",
                "DATARESULTADO": "data_res",
                "REFAZER": "refazer",
            },
            inplace=True,
        )

        df = df[
            [
                "id_proc",
                "cod_empresa_principal",
                "id_empresa",
                "id_unidade",
                "id_funcionario",
                "id_exame",
                "periodicidade",
                "data_adm",
                "ult_pedido",
                "data_res",
                "refazer",
            ]
        ]

        df.dropna(
            axis=0,
            subset=[
                "id_funcionario",
                "id_exame",
                "id_empresa",
                "id_unidade",
            ],
            inplace=True,
        )

        return df

    @classmethod
    def criar_relatorios(
        cls,
        df: pd.DataFrame,
        nome_empresa: str,
        data_origem: datetime,
        nome_unidade: str | None = None,
        gerar_ppt: bool = True,
        pasta: str = UPLOAD_FOLDER,
        filtro_status: list[str] | None = None,
        filtro_a_vencer: list[int] | None = None,
    ) -> str | None:
        """
        Criar CSV e PPT ZIPADOS para o pedido_proc passado

        Retorna nome do arquivo ZIP
        """
        cols = [
            "cod_empresa",
            "razao_social",
            "cod_unidade",
            "nome_unidade",
            "nome_setor",
            "nome_cargo",
            "cod_funcionario",
            "cpf_funcionario",
            "nome_funcionario",
            "cod_exame",
            "nome_exame",
            "periodicidade",
            "data_adm",
            "ult_pedido",
            "data_res",
            "refazer",
            "status",
            "a_vencer",
        ]

        # checar se pasta existe
        if not os.path.exists(pasta):
            os.mkdir(pasta)

        # START--------------------------------------------------------------------------
        df["status"] = list(
            map(cls.criar_status, df["ult_pedido"], df["data_res"], df["refazer"])
        )

        # calcular dias ate o vencimento
        df["dias_vencer"] = list(map(cls.dias_a_vencer, df["refazer"]))
        # substituir todos menores que 1 por NA
        df.loc[df["dias_vencer"] <= 0, "dias_vencer"] = pd.NA  # type: ignore

        # criar categorias a vencer
        df["a_vencer"] = pd.NA
        df.loc[df["dias_vencer"].between(0, 30, "both"), "a_vencer"] = 30
        df.loc[df["dias_vencer"].between(31, 60, "both"), "a_vencer"] = 60
        df.loc[df["dias_vencer"].between(61, 90, "both"), "a_vencer"] = 90
        df.loc[df["dias_vencer"].between(91, 365, "both"), "a_vencer"] = 365

        if filtro_status:
            df = df[df["status"].isin(filtro_status)]

        if filtro_a_vencer:
            df = df[df["a_vencer"].isin(filtro_a_vencer)]

        if not df.empty:
            nome_arquivos = secure_filename(nome_empresa).replace(".", "_")

            timestamp = int(datetime.now().timestamp())

            nome_csv = f"{pasta}/{nome_arquivos}_{timestamp}.xlsx"
            df2 = df[cols]
            df2.to_excel(nome_csv, index=False, freeze_panes=(1, 0))

            arquivos_zipar = [nome_csv]
            gerar_ppt = False
            # gerar ppt apenas se houver mais de 50 exames para refazer
            if gerar_ppt and len(df["refazer"].dropna()) >= 50:
                nome_ppt = f"{pasta}/{nome_arquivos}_{timestamp}.pptx"
                cls.criar_ppt(
                    df=df,
                    data_origem=data_origem,
                    nome_arquivo=nome_ppt,
                    nome_empresa=nome_empresa,
                    nome_unidade=nome_unidade,
                )
                arquivos_zipar.append(nome_ppt)

            nome_zip = f"{pasta}/{nome_arquivos}_{timestamp}.zip"
            nome_zip = zipar_arquivos(
                caminhos_arquivos=arquivos_zipar, caminho_pasta_zip=nome_zip
            )

            return nome_zip
        else:
            return None

    @classmethod
    def criar_relatorios2(
        cls,
        query: BaseQuery,
        nome_empresa: str,
        nome_unidade: str = None,
        gerar_ppt: bool = False,
        data_origem: datetime | None = None,
        filtro_status: list[str] = None,
        filtro_a_vencer: list[int] = None,
    ):
        if gerar_ppt and not data_origem:
            raise ValueError("data_origem é obrigatório se gerar_ppt for True")

        df = pd.read_sql(query.statement, con=db.session.bind)

        if df.empty:
            return None

        df = cls.__tratar_df(df=df)

        if filtro_status:
            df = df[df["id_status"].isin(filtro_status)]

        if filtro_a_vencer:
            if 0 in filtro_a_vencer:
                # substituir zero por pd.NA
                filtro_a_vencer[filtro_a_vencer.index(0)] = pd.NA

            df = df[df["a_vencer"].isin(filtro_a_vencer)]

        if df.empty:
            return None

        # define path for files
        nome = secure_filename(nome_empresa).replace(".", "_").upper()
        timestamp = int(datetime.now().timestamp())
        file_path = os.path.join(UPLOAD_FOLDER, f"{nome}_{timestamp}")

        nome_excel = f"{file_path}.xlsx"
        cls.__gerar_excel(df=df, file_path=nome_excel)

        arquivos = {"excel": nome_excel}

        if gerar_ppt and len(df["refazer"].dropna()) >= cls.PPT_TRIGGER:
            nome_ppt = f"{file_path}.pptx"
            cls.criar_ppt(
                df=df,
                data_origem=data_origem,
                nome_arquivo=nome_ppt,
                nome_empresa=nome_empresa,
                nome_unidade=nome_unidade,
            )
            arquivos["ppt"] = nome_ppt

        return arquivos

    @classmethod
    def __tratar_df(cls, df: pd.DataFrame):
        df = df.copy()

        df["id_status"] = list(
            map(cls.__get_id_status, df["ult_pedido"], df["data_res"], df["refazer"])
        )

        df["status"] = list(map(cls.STATUS_EXAMES.get, df["id_status"]))

        df["dias_vencer"] = list(map(cls.dias_a_vencer, df["refazer"]))
        df.loc[df["dias_vencer"] <= 0, "dias_vencer"] = pd.NA

        df = cls.__criar_cat_a_vencer(df=df)

        return df

    @staticmethod
    def __get_id_status(
        data_ult_pedido: datetime, data_res: datetime, data_refazer: datetime
    ):
        TODAY = datetime.now().date()

        # se nunca foi realizado
        if not data_ult_pedido:
            return 5

        # se ultimo pedido nao foi realizado
        if not data_res:
            return 4

        # se data ja passou ou vence hoje
        if data_refazer <= TODAY:
            return 3
        # se ainda vai vencer ate em 365 dias
        elif data_refazer > TODAY and (data_refazer - TODAY).days <= 365:
            return 2
        # se vai vencer em mais de 365 dias
        elif data_refazer > TODAY and (data_refazer - TODAY).days > 365:
            return 1

    @staticmethod
    def __criar_cat_a_vencer(df: pd.DataFrame):
        df = df.copy()
        df["a_vencer"] = pd.NA
        df.loc[df["dias_vencer"].between(0, 30, "both"), "a_vencer"] = 30
        df.loc[df["dias_vencer"].between(31, 60, "both"), "a_vencer"] = 60
        df.loc[df["dias_vencer"].between(61, 90, "both"), "a_vencer"] = 90
        df.loc[df["dias_vencer"].between(91, 365, "both"), "a_vencer"] = 365
        return df

    @staticmethod
    def __gerar_excel(df: pd.DataFrame, file_path: str):
        COLS_EXCEL = {
            "razao_social": "Empresa",
            "nome_unidade": "Unidade",
            "nome_setor": "Setor",
            "nome_cargo": "Cargo",
            "cpf_funcionario": "CPF Funcionario",
            "nome_funcionario": "Nome Funcionario",
            "nome_exame": "Exame",
            "ult_pedido": "Data Último Exame",
            "periodicidade": "Periodicidade",
            "refazer": "Data Vencimento",
            "status": "Status Exame",
            "cat_vencer": "Vencera em ate",
        }

        df = df.copy()

        df["cat_vencer"] = df["a_vencer"].apply(
            lambda x: f"{x} dias" if x is not pd.NA else None
        )

        # format dates
        for col in ["ult_pedido", "refazer"]:
            df[col] = pd.to_datetime(df[col], dayfirst=True)
            df[col] = df[col].dt.strftime("%d/%m/%Y")

        df_excel: pd.DataFrame = df[COLS_EXCEL.keys()]
        df_excel = df_excel.rename(columns=COLS_EXCEL)
        df_excel.to_excel(file_path, index=False, freeze_panes=(1, 0))

        return None

    @staticmethod
    def criar_status(
        data_ult_pedido: datetime, data_res: datetime, data_refazer: datetime
    ):
        today = datetime.now().date()
        if data_ult_pedido and data_res and data_refazer:
            # se data ja passou ou vence hoje
            if data_refazer <= today:
                return "Vencido"
            # se ainda vai vencer ate em 365 dias
            elif data_refazer > today and (data_refazer - today).days <= 365:
                return "A vencer"
            # se vai vencer em mais de 365 dias
            elif data_refazer > today and (data_refazer - today).days > 365:
                return "Em dia"
        # se nao tem data de resultado
        elif data_ult_pedido and not data_res:
            return "Pendente"
        # se nao tem ultimo pedido
        elif not data_ult_pedido:
            return "Sem histórico"
        else:
            return None

    @staticmethod
    def dias_a_vencer(data: datetime):
        if data:
            dias = (data - datetime.now().date()).days
            if dias > 0:
                return int(dias)
            else:
                return pd.NA
        else:
            return pd.NA

    @classmethod
    def criar_ppt(
        cls,
        df: pd.DataFrame,
        nome_arquivo: str,
        nome_empresa: str,
        data_origem: datetime,
        nome_unidade: str | None = None,
    ):
        # CRIAR PPT--------------------------------------------------------------------
        # para ver os nomes dos shapes
        # for shape in slide.shapes:
        #     print(shape.name)

        df = df.copy()
        # instanciar apresentacao
        presentation = Presentation(cls.BASE_PPT)

        # SLIDE 1 - titulo--------------------------------------------------
        slide = presentation.slides[0]

        df["refazer"] = pd.to_datetime(df["refazer"])

        inicio = df["refazer"].min(skipna=True).strftime("%d/%m/%Y")
        fim = df["refazer"].max(skipna=True).strftime("%d/%m/%Y")
        slide.shapes[4].text_frame.paragraphs[0].runs[-1].text = f"{inicio} - {fim}"

        # nome empresa
        slide.shapes[5].text_frame.paragraphs[0].runs[-1].text = nome_empresa

        # nome unidade
        if nome_unidade:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = nome_unidade
        else:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = "Todas"

        # SLIDE 2 - cards --------------------------------------------------------
        slide = presentation.slides[1]
        shapes = slide.shapes

        # total exames clínicos
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df.query('nome_exame == "Exame Clínico"')))

        # total funcionarios
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(
            len(df["id_funcionario"].drop_duplicates())
        )

        # media exames/funcionario
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            len(df) / len(df["id_funcionario"].drop_duplicates())
        )

        # total exames complementares
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df.query('nome_exame != "Exame Clínico"')))

        # SLIDE 3 exames por mesAno ----------------------------------------------------
        slide = presentation.slides[2]
        shapes = slide.shapes

        df["mesAno"] = df["refazer"].dt.strftime("%m/%Y")

        # qtds = df[["mesAno", "id_exame"]].groupby(by="mesAno").count()

        df["mes"] = df["refazer"].dt.month
        df["ano"] = df["refazer"].dt.year

        dados = (
            df[["ano", "mes", "mesAno", "status", "id_exame"]]
            .pivot_table(
                values="id_exame",
                index=["ano", "mes", "mesAno"],  # gera MultiIndex
                columns="status",
                aggfunc="count",
            )
            .fillna(0)
        )
        # resetar pasa single Index de MesAno
        dados.reset_index(level=["ano", "mes"], drop=True, inplace=True)
        # organizar colunas para manter a ordem no grafico
        dados.loc["Total"] = dados.sum(numeric_only=True)
        dados = dados.sort_values(by="Total", axis=1, ascending=False)
        dados.drop(labels=["Total"], axis=0, inplace=True)

        cls.editar_grafico(  # barras empilhadas
            chart=slide.shapes[2].chart,
            categorias=list(dados.index),
            series={col: list(dados[col].values) for col in dados.columns},
        )

        # SLIDE 4 - exames por status ------------------------------------------------------
        slide = presentation.slides[3]
        shapes = slide.shapes

        # qtds = df[["status", "id_exame"]].groupby(by="status").count()

        # Exames em dia
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df.query('status == "Em dia"')))

        # Exames a vencer
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df.query('status == "A vencer"')))

        # Exames Vencidos
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df.query('status == "Vencido"')))

        # Sem histórico de Exames
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df.query('status == "Sem histórico"')))

        df['exame_complementar'] = 'Exame Complementar'
        df.loc[df['nome_exame'] == 'Exame Clínico', 'exame_complementar'] = 'Exame Clínico'

        dados = df[df['exame_complementar'].isin(['Exame Complementar', 'Exame Clínico'])]

        print(dados.index)

        # dados = df[["status", "id_exame"]].groupby("status").count()
        # dados["id_exame"] = dados["id_exame"] / dados["id_exame"].sum()
        # dados = dados.sort_values("id_exame", axis=0, ascending=False)

        cls.editar_grafico(
            chart=slide.shapes[7].chart,
            categorias=list(dados.index),
            series={"qtd": list(dados["id_exame"].values)},
        )

        # SLIDES DE TABELA 5-9 ------------------------------------------------------
        infos_slides_tabelas = [
            (4, "nome_unidade"),
            (5, "nome_setor"),
            (6, "nome_cargo"),
            (7, "nome_funcionario"),
            (8, "nome_exame"),
        ]

        for num_slide, nome_col in infos_slides_tabelas:
            slide = presentation.slides[num_slide]
            shapes = slide.shapes

            dados = (
                df[[nome_col, "a_vencer", "id_exame"]]
                .pivot_table(
                    values="id_exame",
                    index=nome_col,
                    columns="a_vencer",
                    aggfunc="count",
                )
                .fillna(0)
            )
            dados["Total"] = dados.sum(axis=1, numeric_only=True)
            dados = dados.sort_values(by="Total", axis=0, ascending=False)
            dados.reset_index(inplace=True)

            # checar e organizar as colunas para a tabela do ppt
            for col in [30, 60, 90, 365, "Total"]:
                if col in dados.columns:
                    dados[col] = dados[col].astype(int)
                else:
                    dados[col] = 0

            dados = dados[[nome_col, 30, 60, 90, 365, "Total"]]

            cls.editar_tabela(shape=shapes[3].table, df=dados[:10])

        # SLIDE 10 referecias ------------------------------------------------------
        slide = presentation.slides[9]
        shapes = slide.shapes

        # data de origem dos dados
        slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = data_origem.strftime(
            "%d/%m/%Y"
        )

        presentation.save(nome_arquivo)
        return None

    @staticmethod
    def editar_grafico(chart, categorias: list, series: dict[list]):
        chart_data = ChartData()
        chart_data.categories = categorias
        if series:
            for nome_serie, vals in series.items():
                chart_data.add_series(nome_serie, vals)
        else:
            chart_data.add_series("Sem dados", [0])
        chart.replace_data(chart_data)
        return None

    @staticmethod
    def editar_tabela(shape: object, df: pd.DataFrame, font_size: int = 14):
        for coluna in range(len(df.columns)):
            # pular duas primeiras celulas de cada coluna
            dados_coluna = [None, None] + list(df[df.columns[coluna]].values)
            for linha in range(2, len(dados_coluna)):
                shape.cell(linha, coluna).text_frame.paragraphs[0].text = str(
                    dados_coluna[linha]
                )
                shape.cell(linha, coluna).text_frame.paragraphs[0].runs[
                    0
                ].font.size = Pt(font_size)
        return None

    @classmethod
    def rotina_conv_exames(
        cls,
        id_proc: int,
        corpo_email: str,
        id_unidade: int | None = None,
        testando: bool = True,
        gerar_ppt: bool = True,
        filtro_status: list[str] | None = None,
        filtro_a_vencer: list[int] | None = None,
        tentativas_email: int = 3,
        intervalo_tentativas: int = 30,
    ) -> dict[str, any]:
        """
        Cria os relatorios de Convocacao de Exames e envia para \
        o email da Empresa/Unidade

        Gera ppt apenas se a qtd de linhas for maior que 50

        tentativas_email: numero de tentativas ao enviar o email
        intervalo_tentativas: qtd de segundos entre cada tentaviva

        Returns:
            dict[str, any]: {
                'cod_empresa_principal': int,
                'nome_empresa_principal': str,
                'id_proc': int,
                'id_empresa': int,
                'razao_social': str,
                'id_unidade': int | None,
                'nome_unidade': str | None,
                'data_pedido_proc': datetime.date,
                'emails': str | None,
                'nome_arquivo': str,
                'status': str,
                'erro': str | None,
                'df_len': int,
                'tempo_execucao': int
            }
        """
        start = time.time()

        pedido_proc: PedidoProcessamento = PedidoProcessamento.query.get(id_proc)

        empresa_principal: EmpresaPrincipal = EmpresaPrincipal.query.get(
            pedido_proc.cod_empresa_principal
        )
        empresa: Empresa = Empresa.query.get(pedido_proc.id_empresa)

        infos: dict[str, any] = {
            "cod_empresa_principal": empresa_principal.cod,
            "nome_empresa_principal": empresa_principal.nome,
            "id_proc": pedido_proc.id_proc,
            "id_empresa": empresa.id_empresa,
            "razao_social": empresa.razao_social,
            "id_unidade": None,
            "nome_unidade": None,
            "data_pedido_proc": pedido_proc.data_criacao,
            "emails": empresa.conv_exames_emails,
            "nome_arquivo": None,
            "status": None,
            "erro": None,
            "df_len": None,
            "tempo_execucao": None,
        }

        filtros_query = [(ConvExames.id_proc == pedido_proc.id_proc)]

        if id_unidade:
            unidade: Unidade = Unidade.query.get(id_unidade)

            infos["id_unidade"] = unidade.id_unidade
            infos["nome_unidade"] = unidade.nome_unidade
            infos["emails"] = unidade.conv_exames_emails

            filtros_query.append((ConvExames.id_unidade == infos["id_unidade"]))

        query_conv_exames = (
            db.session.query(ConvExames, Empresa, Unidade, Exame, Funcionario)
            .filter(*filtros_query)
            .outerjoin(Empresa, Empresa.id_empresa == ConvExames.id_empresa)
            .outerjoin(Unidade, Unidade.id_unidade == ConvExames.id_unidade)
            .outerjoin(Exame, Exame.id_exame == ConvExames.id_exame)
            .outerjoin(
                Funcionario, Funcionario.id_funcionario == ConvExames.id_funcionario
            )
        )

        df_conv_exames = pd.read_sql_query(
            sql=query_conv_exames.statement, con=db.session.bind
        )
        infos["df_len"] = len(df_conv_exames)

        if df_conv_exames.empty:
            infos["status"] = "Query ConvExames vazia"
            infos["tempo_execucao"] = int(time.time() - start)
            return infos

        if not infos["emails"]:
            infos["status"] = "Email vazio"
            infos["tempo_execucao"] = int(time.time() - start)
            return infos

        infos["nome_arquivo"] = cls.criar_relatorios(
            df=df_conv_exames,
            nome_empresa=infos["razao_social"],
            nome_unidade=infos["nome_unidade"],
            data_origem=pedido_proc.data_criacao,
            gerar_ppt=gerar_ppt,
            filtro_status=filtro_status,
            filtro_a_vencer=filtro_a_vencer,
        )

        if not infos["nome_arquivo"]:
            infos["status"] = "Sem arquivo"
            infos["tempo_execucao"] = int(time.time() - start)
            return infos

        enviar_para: list[str] = infos["emails"].split(";")
        if testando:
            enviar_para: list[str] = ["gabrielsantos@grsnucleo.com.br"]

        assunto: str = f"Convocação de Exames Empresas - {infos['razao_social']}"
        if id_unidade:
            assunto: str = f"Convocação de Exames Unidades - {infos['nome_unidade']}"

        status_email = EmailConnect.send_email(
            to_addr=enviar_para,
            reply_to=[
                "gabrielsantos@grsnucleo.com.br",
                "relacionamento@grsnucleo.com.br",
            ],
            message_subject=assunto,
            message_body=corpo_email,
            message_imgs=[EmailConnect.ASSINATURA_BOT],
            message_attachments=[infos["nome_arquivo"]],
            send_attempts=tentativas_email,
            attempt_delay=intervalo_tentativas,
        )

        infos["status"] = "OK"
        # pegar erro do email mesmo que o status final seja ok
        infos["erro"] = status_email["error"]
        if status_email["sent"] == 0:
            infos["status"] = "Erro ao enviar email"

        infos["tempo_execucao"] = int(time.time() - start)
        return infos
