import os
from datetime import datetime
from typing import Any

import pandas as pd
from flask_sqlalchemy import BaseQuery
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.util import Pt
from sqlalchemy.exc import DatabaseError, SQLAlchemyError
from werkzeug.utils import secure_filename

from src import TIMEZONE_SAO_PAULO, UPLOAD_FOLDER, app
from src.email_connect import EmailConnect
from src.exporta_dados import ExportaDadosWS
from src.extensions import database as db
from src.main.empresa.empresa import Empresa
from src.main.empresa_principal.empresa_principal import EmpresaPrincipal
from src.main.funcionario.funcionario import Funcionario
from src.main.unidade.unidade import Unidade
from src.utils import get_json_configs, zipar_arquivos


class Absenteismo(db.Model):
    __tablename__ = "Absenteismo"

    id = db.Column(db.Integer, primary_key=True)
    cod_empresa_principal = db.Column(
        db.Integer, db.ForeignKey("EmpresaPrincipal.cod"), nullable=False
    )
    id_funcionario = db.Column(
        db.Integer, db.ForeignKey("Funcionario.id_funcionario"), nullable=False
    )
    id_empresa = db.Column(
        db.Integer, db.ForeignKey("Empresa.id_empresa"), nullable=False
    )
    id_unidade = db.Column(
        db.Integer, db.ForeignKey("Unidade.id_unidade"), nullable=False
    )

    # Licenca (SOCIND)
    tipo_licenca = db.Column(db.String(255))
    cod_medico = db.Column(db.Integer)
    nome_medico = db.Column(db.String(255))
    data_ficha = db.Column(db.Date, nullable=False)
    data_inicio_licenca = db.Column(db.Date)
    data_fim_licenca = db.Column(db.Date)
    afast_horas = db.Column(db.Boolean)
    hora_inicio_licenca = db.Column(db.Integer)  # segundos
    hora_fim_licenca = db.Column(db.Integer)  # segundos
    motivo_licenca = db.Column(db.String(255))
    cid_contestado = db.Column(db.String(255))
    cod_cid = db.Column(db.String(255))
    tipo_cid = db.Column(db.String(255))
    solicitante = db.Column(db.String(255))
    data_inclusao_licenca = db.Column(db.Date)

    # Licença Médica - Informações Básicas
    dias_afastado = db.Column(db.Integer)
    periodo_afastado = db.Column(db.Integer)
    abonado = db.Column(db.Boolean)
    cid = db.Column(db.String(255))

    COLUNAS_PLANILHA = [
        "cod_empresa_principal",
        "cod_empresa",
        "razao_social",
        "cod_unidade",
        "nome_unidade",
        "nome_setor",
        "nome_cargo",
        "cod_funcionario",
        "cpf_funcionario",
        "nome_funcionario",
        "situacao",
        "tipo_licenca",
        "motivo_licenca",
        "data_ficha",
        "data_inicio_licenca",
        "data_fim_licenca",
        "dias_afastado",
        "afast_horas",
        "hora_inicio_licenca",
        "hora_fim_licenca",
        "cid",
        "data_inclusao_licenca",
        "abonado",
        "nome_medico",
        "solicitante",
    ]

    BASE_PPT = "src/main/absenteismo/ppt/base_absenteismo_v2.pptx"

    @classmethod
    def buscar_licencas(
        cls,
        cod_emp_princ: int | None = None,
        id_empresa: int | None = None,
        id_unidade: int | None = None,
        data_inicio: datetime | None = None,
        data_fim: datetime | None = None,
    ) -> BaseQuery:
        """
        Args:
            data_inicio (datetime | None, optional): data da ficha
            data_fim (datetime | None, optional): data da ficha
        """

        filtros = []

        if cod_emp_princ:
            filtros.append(cls.cod_empresa_principal == cod_emp_princ)
        if id_empresa:
            filtros.append(cls.id_empresa == id_empresa)
        if id_unidade:
            filtros.append(cls.id_unidade == id_unidade)
        if data_inicio:
            filtros.append(cls.data_ficha >= data_inicio)
        if data_fim:
            filtros.append(cls.data_ficha <= data_fim)

        query = (
            db.session.query(  # type: ignore
                cls,
                EmpresaPrincipal.nome,
                Empresa.cod_empresa,
                Empresa.razao_social,
                Unidade.cod_unidade,
                Unidade.nome_unidade,
                Funcionario.cod_funcionario,
                Funcionario.cpf_funcionario,
                Funcionario.nome_funcionario,
                Funcionario.nome_setor,
                Funcionario.nome_cargo,
                Funcionario.situacao,
            )
            .outerjoin(
                EmpresaPrincipal, cls.cod_empresa_principal == EmpresaPrincipal.cod
            )
            .outerjoin(Empresa, cls.id_empresa == Empresa.id_empresa)
            .outerjoin(Unidade, cls.id_unidade == Unidade.id_unidade)
            .outerjoin(Funcionario, cls.id_funcionario == Funcionario.id_funcionario)
            .filter(*filtros)
        )

        return query

    @classmethod
    def inserir_licenca(
        cls, id_empresa: int, dataInicio: datetime, dataFim: datetime
    ):
        """Invoca os exporta dados Licenca (SOCIND) e Licença Médica -
        Informações Básicas e insere linhas na tabela.

        Remove licencas duplicadas baseando em 'id_funcionario',
        'tipo_licenca', 'data_inicio_licenca', 'data_fim_licenca'

        Args:
            dataInicio - dataFim: data da ficha

        Returns:
            dict: {'status': str, 'qtd': int}
        """
        empresa: Empresa = Empresa.query.get(id_empresa)

        df_socind = cls.carregar_licenca_socind(
            id_empresa=id_empresa, dataInicio=dataInicio, dataFim=dataFim
        )

        df_licenca_med = cls.carregar_licenca_medica(
            id_empresa=id_empresa, dataInicio=dataInicio, dataFim=dataFim
        )

        infos = {
            "status": "OK",
            "qtd": 0,
        }

        if df_socind is None or df_licenca_med is None:
            infos["status"] = "erro nos dfs"
            return infos

        if df_socind.empty or df_licenca_med.empty:
            infos["status"] = "dfs vazios"
            return infos

        df = df_socind.merge(
            df_licenca_med,
            how="left",
            on=["id_funcionario", "data_ficha"],
        )

        df = cls.__get_infos_licenca(df=df)

        df = df[df["id"].isna()]

        if df.empty:
            infos["status"] = "Sem licencas novas"
            return infos

        FINAL_COLS = [
            "id_funcionario",
            "id_empresa",
            "id_unidade",
            "tipo_licenca",
            "cod_medico",
            "nome_medico",
            "data_ficha",
            "data_inicio_licenca",
            "data_fim_licenca",
            "afast_horas",
            "hora_inicio_licenca",
            "hora_fim_licenca",
            "motivo_licenca",
            "cid_contestado",
            "cod_cid",
            "tipo_cid",
            "solicitante",
            "data_inclusao_licenca",
            "dias_afastado",
            "periodo_afastado",
            "abonado",
            "cid",
        ]

        df = df[FINAL_COLS]

        df["cod_empresa_principal"] = empresa.cod_empresa_principal

        try:
            qtd = df.to_sql(
                Absenteismo.__tablename__,
                con=db.session.bind,  # type: ignore
                index=False,
                if_exists="append",
            )
            db.session.commit()  # type: ignore
        except (SQLAlchemyError, DatabaseError) as err:
            db.session.rollback()  # type: ignore
            app.logger.error(msg=err, exc_info=True)
            infos["status"] = "Erro ao inserir na db"
            return infos

        infos["qtd"] = qtd

        return infos

    @classmethod
    def __get_infos_licenca(cls, df: pd.DataFrame):
        df = df.copy()

        query = db.session.query(  # type: ignore
            Absenteismo.id,
            Absenteismo.id_funcionario,
            Absenteismo.tipo_licenca,
            Absenteismo.data_inicio_licenca,
            Absenteismo.data_fim_licenca,
        )
        df_db = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        # NOTE: make sure date cols are of same type in both dfs
        for col in ("data_inicio_licenca", "data_fim_licenca"):
            df[col] = pd.to_datetime(df[col], errors="coerce", dayfirst=True)
            df_db[col] = pd.to_datetime(df_db[col], errors="coerce", dayfirst=True)

        df = df.merge(
            df_db,
            how="left",
            on=[
                "id_funcionario",
                "tipo_licenca",
                "data_inicio_licenca",
                "data_fim_licenca",
            ],
        )

        return df

    @classmethod
    def carregar_licenca_socind(
        cls, id_empresa: int, dataInicio: datetime, dataFim: datetime
    ):
        """Realiza request para o Exporta Dados Licença (SOCIND).
        Se houver dados na Response, trata o df e retorna os dados
        """
        empresa: Empresa = Empresa.query.get(id_empresa)

        credenciais: dict = get_json_configs(
            empresa.empresa_principal.configs_exporta_dados
        )

        par: dict = ExportaDadosWS.licenca_socind(
            cod_empresa_principal=empresa.cod_empresa_principal,
            cod_exporta_dados=credenciais["LICENCA_SOCIND_COD"],
            chave=credenciais["LICENCA_SOCIND_KEY"],
            empresaTrabalho=empresa.cod_empresa,
            dataInicio=dataInicio.strftime("%d/%m/%Y"),
            dataFim=dataFim.strftime("%d/%m/%Y"),
        )

        response: dict = ExportaDadosWS.request_exporta_dados_ws(parametro=par)

        if response["response"].status_code != 200:
            return None

        erro_soc = response["erro_soc"]
        if erro_soc:
            return None

        response_text: str = response["response"].text
        df = ExportaDadosWS.xml_to_dataframe(xml_string=response_text)

        if df.empty:
            return df

        df = cls.__tto_df_socind(df=df)

        df = cls.__get_infos_funcionario(df=df, id_empresa=empresa.id_empresa)

        df.dropna(axis=0, subset="id_funcionario", inplace=True)

        if df.empty:
            return df

        COLS_FINAIS = [
            "id_empresa",
            "id_unidade",
            "id_funcionario",
            "tipo_licenca",
            "cod_medico",
            "nome_medico",
            "data_ficha",
            "data_inicio_licenca",
            "data_fim_licenca",
            "afast_horas",
            "hora_inicio_licenca",
            "hora_fim_licenca",
            "motivo_licenca",
            "cid_contestado",
            "cod_cid",
            "tipo_cid",
            "solicitante",
            "data_inclusao_licenca",
        ]

        df = df[COLS_FINAIS]

        return df

    @staticmethod
    def __tto_df_socind(df: pd.DataFrame):
        df = df.copy()

        COLS = {
            "CODIGOFUNCIONARIO": "cod_funcionario",
            "TIPO_LICENCA": "tipo_licenca",
            "CODIGO_MEDICO": "cod_medico",
            "MEDICO": "nome_medico",
            "DATA_FICHA": "data_ficha",
            "DATA_INICIO_LICENCA": "data_inicio_licenca",
            "DATA_FIM_LICENCAO": "data_fim_licenca",
            "AFASTAMENTO_EM_HORAS": "afast_horas",
            "HORA_INICIO": "hora_inicio_licenca",
            "HORA_FIM": "hora_fim_licenca",
            "MOTIVO_LICENCA": "motivo_licenca",
            "CID_CONTESTADO": "cid_contestado",
            "CODCID": "cod_cid",
            "TIPO_CID": "tipo_cid",
            "SOLICITANTE": "solicitante",
            "DTINCLUSAOLICENCA": "data_inclusao_licenca",
        }

        df = df[list(COLS.keys())]
        df = df.rename(columns=COLS)

        df = df.replace({"": None})

        INT_COLS = (
            "cod_funcionario",
            "afast_horas",
        )

        for col in INT_COLS:
            df[col] = df[col].astype(int)

        DATE_COLS = (
            "data_ficha",
            "data_inicio_licenca",
            "data_fim_licenca",
            "data_inclusao_licenca",
        )

        for col in DATE_COLS:
            df[col] = pd.to_datetime(df[col], dayfirst=True).dt.date

        df["afast_horas"] = df["afast_horas"].replace(to_replace={1: True, 0: False})

        return df

    @classmethod
    def __get_infos_funcionario(cls, df: pd.DataFrame, id_empresa: int):
        df = df.copy()

        query = db.session.query(  # type: ignore
            Funcionario.id_funcionario,
            Funcionario.cod_funcionario,
            Funcionario.id_empresa,
            Funcionario.id_unidade,
        ).filter(Funcionario.id_empresa == id_empresa)
        df_db = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        df = df.merge(df_db, how="left", on="cod_funcionario")

        return df

    @classmethod
    def carregar_licenca_medica(
        cls, id_empresa: int, dataInicio: datetime, dataFim: datetime
    ):
        """Realiza request para o Exporta Dados Licença Médica - Informações Básicas.
        Se houver dados na Response, trata o df e retorna os dados
        """
        empresa: Empresa = Empresa.query.get(id_empresa)

        credenciais: dict = get_json_configs(
            empresa.empresa_principal.configs_exporta_dados
        )

        par: dict = ExportaDadosWS.licenca_medica(
            cod_exporta_dados=credenciais["LICENCA_MED_COD"],
            chave=credenciais["LICENCA_MED_KEY"],
            empresaTrabalho=empresa.cod_empresa,
            dataInicio=dataInicio.strftime("%d/%m/%Y"),
            dataFim=dataFim.strftime("%d/%m/%Y"),
        )

        response: dict = ExportaDadosWS.request_exporta_dados_ws(parametro=par)

        if response["response"].status_code != 200:
            return None

        erro_soc = response["erro_soc"]
        if erro_soc:
            return None

        response_text: str = response["response"].text
        df = ExportaDadosWS.xml_to_dataframe(xml_string=response_text)

        if df.empty:
            return df

        df = cls.__tto_df_licenca_med(df=df)

        df = cls.__get_infos_funcionario(df=df, id_empresa=empresa.id_empresa)

        df.dropna(axis=0, subset="id_funcionario", inplace=True)

        if df.empty:
            return df

        # remover licencas duplicadas
        df["chave"] = df["id_funcionario"].astype(str) + df["data_ficha"].astype(str)
        df.drop_duplicates(subset="chave", inplace=True, ignore_index=True)
        df.drop(columns="chave", inplace=True)

        FINAL_COLS = [
            "id_funcionario",
            "data_ficha",
            "dias_afastado",
            "periodo_afastado",
            "abonado",
            "cid",
        ]

        df = df[FINAL_COLS]

        return df

    @staticmethod
    def __tto_df_licenca_med(df: pd.DataFrame):
        df = df.copy()

        COLS = {
            "FUNCIONARIO": "cod_funcionario",
            "DATAFICHA": "data_ficha",
            "DIASAFASTADOS": "dias_afastado",
            "PEIODOAFASTADO": "periodo_afastado",
            "ABONADO": "abonado",
            "CID": "cid",
        }

        df = df.replace({"": None})

        df = df[list(COLS.keys())]
        df = df.rename(columns=COLS)

        for col in ("cod_funcionario", "dias_afastado", "abonado"):
            df[col] = df[col].astype(int)

        df["data_ficha"] = pd.to_datetime(df["data_ficha"], dayfirst=True).dt.date

        df["abonado"] = df["abonado"].replace(to_replace={1: True, 2: False})

        return df

    @classmethod
    def rotina_absenteismo(
        cls,
        cod_empresa_principal: int,
        id_empresa: int,
        nome_empresa: str,
        data_inicio: datetime,
        data_fim: datetime,
        emails_destinatario: list[str],
        corpo_email: str,
        id_unidade: int | None = None,
        nome_unidade: str | None = None,
        testando: bool = False,
    ) -> dict[str, Any]:
        """
        Gera os relatorios Excel e PPT, zipa e envia por email.

        Gera ppt apenas se o DataFrame tiver mais de 50 linhas

        Returns:
            dict: {status: str, qtd: int}
        """
        query = cls.buscar_licencas(
            cod_emp_princ=cod_empresa_principal,
            id_empresa=id_empresa,
            id_unidade=id_unidade,
            data_inicio=data_inicio,
            data_fim=data_fim,
        )

        df = pd.read_sql(sql=query.statement, con=db.session.bind)  # type: ignore

        infos = {"status": "OK", "qtd": 0}

        if df.empty:
            infos["status"] = "df vazio"
            return infos

        timestamp = int(datetime.now().timestamp())
        nome_entidade = secure_filename(nome_empresa).replace(".", "_")

        if nome_unidade:
            nome_entidade = secure_filename(nome_unidade).replace(".", "_")

        nome_arquivo = f"Absenteismo_{nome_entidade}_{timestamp}"
        caminho_arqvs = os.path.join(UPLOAD_FOLDER, nome_arquivo)

        infos["qtd"] = len(df)

        df["hora_inicio_licenca"] = pd.to_datetime(
            df["hora_inicio_licenca"], unit="s"
        ).dt.strftime("%H:%M")

        df["hora_fim_licenca"] = pd.to_datetime(
            df["hora_fim_licenca"], unit="s"
        ).dt.strftime("%H:%M")

        df2 = df[cls.COLUNAS_PLANILHA]
        nome_excel = f"{caminho_arqvs}.xlsx"
        df2.to_excel(nome_excel, index=False, freeze_panes=(1, 0))
        arquivos_zipar: list[str] = [nome_excel]

        if id_unidade:
            qtd_ativos = (
                db.session.query(Funcionario)  # type: ignore
                .filter(Funcionario.situacao == "Ativo")
                .filter(Funcionario.id_unidade == id_unidade)
                .count()
            )
        else:
            qtd_ativos = (
                db.session.query(Funcionario)  # type: ignore
                .filter(Funcionario.situacao == "Ativo")
                .filter(Funcionario.id_empresa == id_empresa)
                .count()
            )

        if len(df) >= 50:
            nome_ppt = f"{caminho_arqvs}.pptx"
            arquivos_zipar.append(nome_ppt)
            cls.criar_ppt(
                df=df,
                funcionarios_ativos=qtd_ativos,
                nome_arquivo=nome_ppt,
                nome_empresa=nome_empresa,
                nome_unidade=nome_unidade,
            )

        pasta_zip = zipar_arquivos(
            caminhos_arquivos=arquivos_zipar,
            caminho_pasta_zip=f"{caminho_arqvs}.zip",
        )

        if nome_unidade:
            ass = f"Absenteísmo Unidade - {nome_unidade}"
        else:
            ass = f"Absenteísmo Empresa - {nome_empresa}"

        if testando:
            enviar_para = ["gabrielsantos@grsnucleo.com.br"]
        else:
            enviar_para = emails_destinatario

        email_enviado = True
        erro_email = None
        try:
            EmailConnect.send_email(
                to_addr=enviar_para,
                message_subject=ass,
                message_body=corpo_email,
                message_imgs=[EmailConnect.ASSINATURA_BOT],
                message_attachments=[pasta_zip],
                reply_to=["gabrielsantos@grsnucleo.com.br"],
            )
        except Exception as err:
            infos["status"] = f"erro ao enviar email {str(err)}"
            email_enviado = False
            erro_email = type(err).__name__

        log_email = EmailConnect(
            email_to=",".join(enviar_para),
            email_subject=ass,
            attachments=pasta_zip,
            status=email_enviado,
            error=erro_email,
            df_len=len(df),
            email_date=datetime.now(tz=TIMEZONE_SAO_PAULO),
        )
        db.session.add(log_email)  # type: ignore
        db.session.commit()  # type: ignore

        return infos

    @classmethod
    def criar_ppt(
        cls,
        df: pd.DataFrame,
        funcionarios_ativos: int,
        nome_arquivo: str,
        nome_empresa: str,
        nome_unidade: str | None = None,
    ):
        df = df.copy()
        # instanciar apresentacao
        presentation = Presentation(cls.BASE_PPT)

        # SLIDE 1  - titulo--------------------------------------------------
        slide = presentation.slides[0]

        # periodo
        inicio = df["data_ficha"].min().strftime("%d/%m/%Y")
        fim = df["data_ficha"].max().strftime("%d/%m/%Y")
        slide.shapes[4].text_frame.paragraphs[0].runs[-1].text = f"{inicio} - {fim}"

        # nome empresa
        slide.shapes[5].text_frame.paragraphs[0].runs[-1].text = f"{nome_empresa}"

        # nome unidade
        if nome_unidade:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = f"{nome_unidade}"
        else:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = "Todas"

        # SLIDE 2 - cards --------------------------------------------------------
        slide = presentation.slides[1]
        shapes = slide.shapes

        # total funcionarios ativos
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = str(
            funcionarios_ativos
        )

        # total de licencas
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df))

        # total dias afastados
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(
            df["dias_afastado"].sum()
        )

        # total func afastados
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(
            len(df["id_funcionario"].drop_duplicates())
        )

        # licencas > 15 dias
        shapes[7].shapes[1].text_frame.paragraphs[0].runs[0].text = str(
            len(df[df["dias_afastado"] > 15])
        )

        # taxa absenteismo: horas perdidas / horas previstas
        # shapes[14].text_frame.paragraphs[0].runs[0].text = str()

        # media licencas/funcionario (IFL qtd de licencas / qtd funcionarios ativos)
        shapes[8].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            len(df) / len(df["id_funcionario"].drop_duplicates())
        )

        # media dias afast/funcionario (IG dias afastados / qtd funcionarios ativos)
        shapes[9].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            df["dias_afastado"].sum() / len(df["id_funcionario"].drop_duplicates())
        )

        # media dias por licenca (IDA dias afastados / qtd de funcionarios afastados)
        shapes[10].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            df["dias_afastado"].mean(skipna=True)
        )

        # SLIDE 3 ----------------------------------------------------
        slide = presentation.slides[2]
        shapes = slide.shapes

        df["data_ficha"] = pd.to_datetime(df["data_ficha"])
        df["mes"] = df["data_ficha"].dt.month
        df["ano"] = df["data_ficha"].dt.year
        df["mesAno"] = df["data_ficha"].dt.strftime("%m/%Y")

        qtds = (
            df[["id_funcionario"]]
            .groupby(by=[df["ano"], df["mes"], df["mesAno"]])
            .count()
        )

        qtds.reset_index(inplace=True, level=["mes", "ano"], drop=True)

        # media licencas/mes
        media = qtds.mean(skipna=True).values[0]
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # min licencas/mes
        minimo = qtds.min(skipna=True).values[0]
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # max licencas/mes
        maximo = qtds.max(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        cls.editar_grafico(
            chart=slide.shapes[6].chart,
            categorias=list(qtds.index),
            series={"qtd": list(qtds["id_funcionario"].values)},
        )

        # SLIDE 4 ----------------------------------------------------
        slide = presentation.slides[3]
        shapes = slide.shapes

        qtds = (
            df[["id_funcionario"]]
            .drop_duplicates(subset="id_funcionario")
            .groupby(by=[df["ano"], df["mes"], df["mesAno"]])
            .count()
        )

        qtds.reset_index(inplace=True, level=["mes", "ano"], drop=True)

        # media  funcionarios/mes
        media = qtds.mean(skipna=True).values[0]
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = "{:.2f}".format(
            media
        )

        # min funcionarios/mes
        minimo = qtds.min(skipna=True).values[0]
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # max funcionarios/mes
        maximo = qtds.max(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        cls.editar_grafico(
            chart=slide.shapes[6].chart,
            categorias=list(qtds.index),
            series={"qtd": list(qtds["id_funcionario"].values)},
        )

        # SLIDE 5 ----------------------------------------------------
        slide = presentation.slides[4]
        shapes = slide.shapes

        # qtd licencas por dia da semana
        dias = {0: "Seg", 1: "Ter", 2: "Qua", 3: "Qui", 4: "Sex", 5: "Sáb", 6: "Dom"}
        df["dia_da_sem"] = df["data_ficha"].dt.dayofweek
        df = df.replace(to_replace={"dia_da_sem": dias})

        dados = df["dia_da_sem"].value_counts()
        indice = {"Dom": 0, "Seg": 1, "Ter": 2, "Qua": 3, "Qui": 4, "Sex": 5, "Sáb": 6}
        dados = dados.sort_index(key=lambda x: x.map(indice))

        cls.editar_grafico(
            chart=shapes[3].chart,
            categorias=list(dados.index),
            series={"qtd": list(dados.values)},
        )

        # qtd licencas por cid
        # remover cids vazios para evitar contagem de valores NaN
        dados = df.dropna(subset=["cid"])
        dados["grupo_cid"] = dados["cid"].astype(str).str[0].str.upper()
        dados = dados["grupo_cid"].value_counts().sort_values(ascending=False)[:10]
        cls.editar_grafico(
            chart=shapes[4].chart,
            categorias=list(dados.index),
            series={"qtd": list(dados.values)},
        )

        # qtd licencas por dias
        dados = df["dias_afastado"].value_counts().sort_values(ascending=False)[:10]
        cls.editar_grafico(
            chart=shapes[5].chart,
            categorias=list(dados.index),
            series={"qtd": list(dados.values)},
        )

        # licencas abonadas
        dados = df["abonado"].value_counts().reset_index()
        dados = dados.replace({True: "Sim", False: "Não"})
        dados["abonado"] = dados["abonado"] / len(df)  # calc percentuais
        dados.sort_values(by="abonado", ascending=False, inplace=True)
        cls.editar_grafico(
            chart=shapes[6].chart,
            categorias=list(dados["index"].values),
            series={"Abonado": list(dados["abonado"].values)},
        )

        # # SLIDE 6-8 RANKINGS  ----------------------------------------
        infos_slides_tabelas = [
            (5, "nome_unidade"),
            (6, "nome_setor"),
            (7, "nome_cargo"),
            (8, "nome_funcionario"),
        ]

        # adicionar col para usar na soma de licencas
        df["qtd_licencas"] = 1

        for num_slide, nome_col in infos_slides_tabelas:
            slide = presentation.slides[num_slide]
            shapes = slide.shapes

            dados = (
                df[[nome_col, "qtd_licencas", "dias_afastado"]]
                .groupby(nome_col)
                .sum()
                .fillna(0)
            )
            dados = dados.reset_index()
            dados = dados.sort_values("dias_afastado", ascending=False)

            cls.editar_tabela(shape=shapes[3].table, df=dados[:10])

        # SLIDE 10 referecias ------------------------------------------------------
        slide = presentation.slides[9]
        shapes = slide.shapes

        # relatorio gerado em
        slide.shapes[6].text_frame.paragraphs[0].runs[
            -1
        ].text = datetime.now().strftime("%d/%m/%Y")

        presentation.save(nome_arquivo)
        return None

    @staticmethod
    def editar_grafico(chart, categorias: list, series: dict[str, list[int]]):
        chart_data = ChartData()
        if series and categorias:
            chart_data.categories = categorias
            for nome_serie, vals in series.items():
                chart_data.add_series(nome_serie, vals)
        else:
            chart_data.categories = ["Sem dados"]
            chart_data.add_series("Sem dados", [0])
        chart.replace_data(chart_data)
        return None

    @staticmethod
    def editar_tabela(shape: object, df: pd.DataFrame, font_size: int = 14):
        for coluna in range(len(df.columns)):
            # pular duas primeiras celulas de cada coluna
            dados_coluna = [None, None] + list(df[df.columns[coluna]].values)
            for linha in range(2, len(dados_coluna)):
                shape.cell(linha, coluna).text_frame.paragraphs[0].text = str(  # type: ignore
                    dados_coluna[linha]
                )
                shape.cell(linha, coluna).text_frame.paragraphs[0].runs[  # type: ignore
                    0
                ].font.size = Pt(font_size)
        return None
