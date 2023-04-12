from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.util import Pt
from werkzeug.utils import secure_filename

from src import TIMEZONE_SAO_PAULO, UPLOAD_FOLDER, database
from src.email_connect import EmailConnect
from src.exporta_dados import ExportaDadosWS
from src.main.empresa.empresa import Empresa
from src.main.empresa_principal.empresa_principal import EmpresaPrincipal
from src.main.funcionario.funcionario import Funcionario
from src.main.unidade.unidade import Unidade
from src.utils import get_json_configs, zipar_arquivos


class Licenca(database.Model):
    __tablename__ = 'Licenca'
    id = database.Column(database.Integer, primary_key=True)
    cod_empresa_principal = database.Column(database.Integer, database.ForeignKey('EmpresaPrincipal.cod'), nullable=False)
    id_funcionario = database.Column(database.Integer, database.ForeignKey('Funcionario.id_funcionario'), nullable=False)
    id_empresa = database.Column(database.Integer, database.ForeignKey('Empresa.id_empresa'), nullable=False)
    id_unidade = database.Column(database.Integer, database.ForeignKey('Unidade.id_unidade'), nullable=False)
    # Licenca (SOCIND)
    tipo_licenca = database.Column(database.String(255))
    cod_medico = database.Column(database.Integer)
    nome_medico = database.Column(database.String(255))
    data_ficha = database.Column(database.Date, nullable=False)
    data_inicio_licenca = database.Column(database.Date)
    data_fim_licenca = database.Column(database.Date)
    afast_horas = database.Column(database.Boolean)
    hora_inicio_licenca = database.Column(database.Integer) # segundos
    hora_fim_licenca = database.Column(database.Integer) # segundos
    motivo_licenca = database.Column(database.String(255))
    cid_contestado = database.Column(database.String(255))
    cod_cid  = database.Column(database.String(255))
    tipo_cid = database.Column(database.String(255))
    solicitante = database.Column(database.String(255))
    data_inclusao_licenca = database.Column(database.Date)
    # Licença Médica - Informações Básicas
    dias_afastado = database.Column(database.Integer)
    periodo_afastado = database.Column(database.Integer)
    abonado = database.Column(database.Boolean)
    cid = database.Column(database.String(255))

    COLUNAS_PLANILHA = [
        'cod_empresa_principal',
        'cod_empresa',
        'razao_social',
        'cod_unidade',
        'nome_unidade',
        'nome_setor',
        'nome_cargo',
        'cod_funcionario',
        'cpf_funcionario',
        'nome_funcionario',
        'situacao',
        'tipo_licenca',
        'motivo_licenca',
        'data_ficha',
        'data_inicio_licenca',
        'data_fim_licenca',
        'dias_afastado',
        'afast_horas',
        'hora_inicio_licenca',
        'hora_fim_licenca',
        'cid',
        'data_inclusao_licenca',
        'abonado',
        'nome_medico',
        'solicitante'
    ]

    BASE_PPT = 'src/modules/absenteismo/ppt/base_absenteismo_v2.pptx'


    @classmethod
    def buscar_licencas(
        self,
        cod_empresa_principal: int | None = None,
        id_empresa: int = None,
        id_unidade: int = None,
        data_inicio: datetime = None,
        data_fim: datetime = None
    ):

        filtros = []

        if cod_empresa_principal:
            filtros.append(self.cod_empresa_principal == cod_empresa_principal)
        if id_empresa:
            filtros.append(self.id_empresa == id_empresa)
        if id_unidade:
            filtros.append(self.id_unidade == id_unidade)
        if data_inicio:
            filtros.append(self.data_ficha >= data_inicio)
        if data_fim:
            filtros.append(self.data_ficha <= data_fim)
        
        query = (
            database.session.query(
                self,
                EmpresaPrincipal.nome,
                Empresa.cod_empresa,
                Empresa.razao_social,
                Unidade.cod_unidade,
                Unidade.nome_unidade,
                Funcionario.cod_funcionario,
                Funcionario.cpf_funcionario,
                Funcionario.nome_funcionario,
                Funcionario.nome_setor,
                Funcionario.nome_cargo,
                Funcionario.situacao
            )
            .outerjoin(EmpresaPrincipal, self.cod_empresa_principal == EmpresaPrincipal.cod)
            .outerjoin(Empresa, self.id_empresa == Empresa.id_empresa)
            .outerjoin(Unidade, self.id_unidade == Unidade.id_unidade)
            .outerjoin(Funcionario, self.id_funcionario == Funcionario.id_funcionario)
            .filter(*filtros)
        )

        return query

    # CARREGAMENTO DADOS ----------------------------------------------------------
    @classmethod
    def inserir_licenca(
        self,
        id_empresa: int,
        dataInicio: datetime,
        dataFim: datetime
    ) -> dict[dict]:
        """Invoca os exporta dados Licenca (SOCIND) e Licença Médica \
        - Informações Básicas e insere linhas na tabela.

        Remove licencas duplicadas baseando em 'id_funcionario', \
        'tipo_licenca', 'data_inicio_licenca', 'data_fim_licenca'

        Args:
            dataInicio - dataFim: data da ficha

        Returns:
        Retorna os dicionarios das duas funcoes usadas
            dict[dict]: {
                'geral': geral,
                'socind': licenca_socind,
                'licenca_med': licenca_med
            } 
        """

        empresa: Empresa = Empresa.query.get(id_empresa)
        empresa_principal: EmpresaPrincipal = EmpresaPrincipal.query.get(empresa.cod_empresa_principal)

        licenca_socind: dict = self.carregar_licenca_socind(
            id_empresa = id_empresa,
            dataInicio = dataInicio,
            dataFim = dataFim
        )

        licenca_med: dict = self.carregar_licenca_medica(
            id_empresa = id_empresa,
            dataInicio = dataInicio,
            dataFim = dataFim
        )

        df_socind: pd.DataFrame = licenca_socind['df']
        df_licenca_med: pd.DataFrame = licenca_med['df']

        if not df_socind.empty and not df_licenca_med.empty:
            df = df_socind.merge(
                df_licenca_med,
                how='left',
                on=['id_funcionario', 'data_ficha'],
            )

            # remover duplicados
            df_db = pd.read_sql(
                sql=(
                    database.session.query(
                        Licenca.id,
                        Licenca.id_funcionario,
                        Licenca.tipo_licenca,
                        Licenca.data_inicio_licenca,
                        Licenca.data_fim_licenca
                    )
                ).statement,
                con=database.session.bind
            )

            df = df.merge(
                df_db,
                how='left',
                on=[
                    'id_funcionario',
                    'tipo_licenca',
                    'data_inicio_licenca',
                    'data_fim_licenca'
                ]
            )

            df = df[df['id'].isna()]

            df = df[[
                'id_funcionario',
                'id_empresa',
                'id_unidade',
                'tipo_licenca',
                'cod_medico',
                'nome_medico',
                'data_ficha',
                'data_inicio_licenca',
                'data_fim_licenca',
                'afast_horas',
                'hora_inicio_licenca',
                'hora_fim_licenca',
                'motivo_licenca',
                'cid_contestado',
                'cod_cid',
                'tipo_cid',
                'solicitante',
                'data_inclusao_licenca',
                'dias_afastado',
                'periodo_afastado',
                'abonado',
                'cid'
            ]]

            df['cod_empresa_principal'] = empresa.cod_empresa_principal

            qtd = df.to_sql(Licenca.__tablename__, con=database.session.bind, index=False, if_exists='append')
            database.session.commit()

            geral = {
                'cod_empresa_principal': empresa_principal.cod,
                'nome_empresa_principal': empresa_principal.nome,
                'id_empresa': empresa.id_empresa,
                'nome_empresa': empresa.razao_social,
                'status': 'ok',
                'df_licenca_med': len(df_licenca_med),
                'df_socind': len(df_socind),
                'qtd': qtd
            }
            licenca_socind['df'] = len(licenca_socind['df'])
            licenca_med['df'] = len(licenca_med['df'])
            return {
                'geral': geral,
                'socind': licenca_socind,
                'licenca_med': licenca_med
            }
        else:
            geral = {
                'cod_empresa_principal': empresa_principal.cod,
                'nome_empresa_principal': empresa_principal.nome,
                'id_empresa': empresa.id_empresa,
                'nome_empresa': empresa.razao_social,
                'status': 'erro/vazio',
                'df_licenca_med': len(df_licenca_med),
                'df_socind': len(df_socind),
                'qtd': 0
            }
            licenca_socind['df'] = len(licenca_socind['df'])
            licenca_med['df'] = len(licenca_med['df'])
            return {
                'geral': geral,
                'socind': licenca_socind,
                'licenca_med': licenca_med
            }

    @classmethod
    def carregar_licenca_socind(
        self,
        id_empresa: int,
        dataInicio: datetime,
        dataFim: datetime
    ):
        """Realiza request para o Exporta Dados Licença (SOCIND). \
        Se houver dados na Response, trata o df e retorna os dados

        Returns:
            dict[str, any]: {
                'cod_empresa_principal': int,
                'nome_empresa_principal': str,
                'id_empresa': int,
                'nome_empresa': str,
                'status': str,
                'erro_soc': bool,
                'erro_request': bool,
                'df': pd.DataFrame
            }
        """
        

        empresa: Empresa = Empresa.query.get(id_empresa)
        empresa_principal: EmpresaPrincipal = EmpresaPrincipal.query.get(empresa.cod_empresa_principal)
        credenciais: dict = get_json_configs(empresa_principal.configs_exporta_dados)

        par: dict = ExportaDadosWS.licenca_socind(
            cod_empresa_principal = empresa.cod_empresa_principal,
            cod_exporta_dados = credenciais['LICENCA_SOCIND_COD'],
            chave = credenciais['LICENCA_SOCIND_KEY'],
            empresaTrabalho = empresa.cod_empresa,
            dataInicio = dataInicio.strftime('%d/%m/%Y'),
            dataFim = dataFim.strftime('%d/%m/%Y')
        )

        response: dict = ExportaDadosWS.request_exporta_dados_ws(
            parametro=par,
            id_empresa=empresa.id_empresa,
            obs='Licenca (SOCIND)'
        )

        if response['response'].status_code == 200:
            if not response['erro_soc']:
                response_text: str = response['response'].text
                df = ExportaDadosWS.xml_to_dataframe(xml_string=response_text)
                if not df.empty:
                    df = df.replace({'': None})

                    df['CODIGOFUNCIONARIO'] = df['CODIGOFUNCIONARIO'].astype(int)

                    # buscar ids
                    df_database = pd.read_sql(
                        sql= (
                            database.session.query(
                                Funcionario.id_funcionario,
                                Funcionario.cod_funcionario,
                                Funcionario.id_empresa,
                                Funcionario.id_unidade,
                            )
                            .filter(Funcionario.cod_empresa_principal == empresa.cod_empresa_principal)
                            .filter(Funcionario.id_empresa == empresa.id_empresa)
                        ).statement,
                        con=database.session.bind
                    )

                    df = pd.merge(
                        df,
                        df_database,
                        how='left',
                        left_on='CODIGOFUNCIONARIO',
                        right_on='cod_funcionario'
                    )

                    df.dropna(axis=0, subset='id_funcionario', inplace=True)

                    if not df.empty:
                        # tratar colunas
                        for col in ['DATA_FICHA', 'DATA_INICIO_LICENCA', 'DATA_FIM_LICENCAO', 'DTINCLUSAOLICENCA']:
                            df[col] = pd.to_datetime(df[col], dayfirst=True).dt.date
                        
                        df['AFASTAMENTO_EM_HORAS'] = df['AFASTAMENTO_EM_HORAS'].astype(int)
                        df['AFASTAMENTO_EM_HORAS'] = df['AFASTAMENTO_EM_HORAS'].replace(to_replace={1: True, 0: False})

                        df.rename(
                            columns={
                                'TIPO_LICENCA': 'tipo_licenca',
                                'CODIGO_MEDICO': 'cod_medico',
                                'MEDICO': 'nome_medico',
                                'DATA_FICHA': 'data_ficha',
                                'DATA_INICIO_LICENCA': 'data_inicio_licenca',
                                'DATA_FIM_LICENCAO': 'data_fim_licenca',
                                'AFASTAMENTO_EM_HORAS': 'afast_horas',
                                'HORA_INICIO': 'hora_inicio_licenca',
                                'HORA_FIM': 'hora_fim_licenca',
                                'MOTIVO_LICENCA': 'motivo_licenca',
                                'CID_CONTESTADO': 'cid_contestado',
                                'CODCID': 'cod_cid',
                                'TIPO_CID': 'tipo_cid',
                                'SOLICITANTE': 'solicitante',
                                'DTINCLUSAOLICENCA': 'data_inclusao_licenca'
                            },
                            inplace=True
                        )

                        df = df[[
                            'id_empresa',
                            'id_unidade',
                            'id_funcionario',
                            'tipo_licenca',
                            'cod_medico',
                            'nome_medico',
                            'data_ficha',
                            'data_inicio_licenca',
                            'data_fim_licenca',
                            'afast_horas',
                            'hora_inicio_licenca',
                            'hora_fim_licenca',
                            'motivo_licenca',
                            'cid_contestado',
                            'cod_cid',
                            'tipo_cid',
                            'solicitante',
                            'data_inclusao_licenca'
                        ]]
                        return {
                            'cod_empresa_principal': empresa_principal.cod,
                            'nome_empresa_principal': empresa_principal.nome,
                            'id_empresa': empresa.id_empresa,
                            'nome_empresa': empresa.razao_social,
                            'status': "ok",
                            'erro_soc': False,
                            'erro_request': False,
                            'df': df
                        }
                    else:
                        return {
                            'cod_empresa_principal': empresa_principal.cod,
                            'nome_empresa_principal': empresa_principal.nome,
                            'id_empresa': empresa.id_empresa,
                            'nome_empresa': empresa.razao_social,
                            'status': "vazio",
                            'erro_soc': False,
                            'erro_request': False,
                            'df': pd.DataFrame()
                        }
                else:
                    return {
                        'cod_empresa_principal': empresa_principal.cod,
                        'nome_empresa_principal': empresa_principal.nome,
                        'id_empresa': empresa.id_empresa,
                        'nome_empresa': empresa.razao_social,
                        'status': "vazio",
                        'erro_soc': False,
                        'erro_request': False,
                        'df': pd.DataFrame()
                    }
            else:
                return {
                    'cod_empresa_principal': empresa_principal.cod,
                    'nome_empresa_principal': empresa_principal.nome,
                    'id_empresa': empresa.id_empresa,
                    'nome_empresa': empresa.razao_social,
                    'status': f"erro soc: {response['msg_erro']}",
                    'erro_soc': True,
                    'erro_request': False,
                    'df': pd.DataFrame()
                }
        else:
            return {
                'cod_empresa_principal': empresa_principal.cod,
                'nome_empresa_principal': empresa_principal.nome,
                'id_empresa': empresa.id_empresa,
                'nome_empresa': empresa.razao_social,
                'status': 'erro no request',
                'erro_soc': False,
                'erro_request': True,
                'df': pd.DataFrame()
            }

    @classmethod
    def carregar_licenca_medica(
        self,
        id_empresa: int,
        dataInicio: datetime,
        dataFim: datetime
    ):
        """Realiza request para o Exporta Dados Licença Médica - Informações Básicas \
        Se houver dados na Response, trata o df e retorna os dados

        Returns:
            dict[str, any]: {
                'cod_empresa_principal': int,
                'nome_empresa_principal': str,
                'id_empresa': int,
                'nome_empresa': str,
                'status': str,
                'erro_soc': bool,
                'erro_request': bool,
                'df': pd.DataFrame
            }
        """
        empresa: Empresa = Empresa.query.get(id_empresa)
        empresa_principal: EmpresaPrincipal = EmpresaPrincipal.query.get(empresa.cod_empresa_principal)
        credenciais: dict = get_json_configs(empresa_principal.configs_exporta_dados)

        par: dict = ExportaDadosWS.licenca_medica(
            cod_exporta_dados = credenciais['LICENCA_MED_COD'],
            chave = credenciais['LICENCA_MED_KEY'],
            empresaTrabalho = empresa.cod_empresa,
            dataInicio = dataInicio.strftime('%d/%m/%Y'),
            dataFim = dataFim.strftime('%d/%m/%Y')
        )

        response: dict = ExportaDadosWS.request_exporta_dados_ws(
            parametro=par,
            id_empresa=empresa.id_empresa,
            obs='Licença Médica - Informações Básicas'
        )

        if response['response'].status_code == 200:
            if not response['erro_soc']:
                response_text: str = response['response'].text
                df = ExportaDadosWS.xml_to_dataframe(xml_string=response_text)
                if not df.empty:
                    df = df.replace({'': None})

                    df['FUNCIONARIO'] = df['FUNCIONARIO'].astype(int)

                    # buscar ids
                    df_database = pd.read_sql(
                        sql= (
                            database.session.query(
                                Funcionario.id_funcionario,
                                Funcionario.cod_funcionario,
                                Funcionario.id_empresa,
                                Funcionario.id_unidade,
                            )
                            .filter(Funcionario.cod_empresa_principal == empresa.cod_empresa_principal)
                            .filter(Funcionario.id_empresa == empresa.id_empresa)
                        ).statement,
                        con=database.session.bind
                    )

                    df = pd.merge(
                        df,
                        df_database,
                        how='left',
                        left_on='FUNCIONARIO',
                        right_on='cod_funcionario'
                    )

                    df.dropna(axis=0, subset='id_funcionario', inplace=True)
                    if not df.empty:
                        # tratar colunas
                        df['DATAFICHA'] = pd.to_datetime(df['DATAFICHA'], dayfirst=True).dt.date
                        df['ABONADO'] = df['ABONADO'].astype(int)
                        df['ABONADO'] = df['ABONADO'].replace(to_replace={1:True, 2: False})

                        df['chave'] = (
                            df['id_funcionario'].astype(str) + 
                            df['DATAFICHA'].astype(str)
                        )

                        df.drop_duplicates(subset='chave', inplace=True, ignore_index=True)

                        df.rename(
                            columns={
                                'DATAFICHA': 'data_ficha',
                                'DIASAFASTADOS': 'dias_afastado',
                                'PEIODOAFASTADO': 'periodo_afastado',
                                'ABONADO': 'abonado',
                                'CID': 'cid'
                            },
                            inplace=True
                        )

                        df = df[[
                            'id_funcionario',
                            'data_ficha',
                            'dias_afastado',
                            'periodo_afastado',
                            'abonado',
                            'cid'
                        ]]

                        return {
                            'cod_empresa_principal': empresa_principal.cod,
                            'nome_empresa_principal': empresa_principal.nome,
                            'id_empresa': empresa.id_empresa,
                            'nome_empresa': empresa.razao_social,
                            'status': 'ok',
                            'erro_soc': False,
                            'erro_request': False,
                            'df': df
                        }
                    else:
                        return {
                            'cod_empresa_principal': empresa_principal.cod,
                            'nome_empresa_principal': empresa_principal.nome,
                            'id_empresa': empresa.id_empresa,
                            'nome_empresa': empresa.razao_social,
                            'status': 'vazio',
                            'erro_soc': False,
                            'erro_request': False,
                            'df': pd.DataFrame()
                        }
                else:
                    return {
                        'cod_empresa_principal': empresa_principal.cod,
                        'nome_empresa_principal': empresa_principal.nome,
                        'id_empresa': empresa.id_empresa,
                        'nome_empresa': empresa.razao_social,
                        'status': 'vazio',
                        'erro_soc': False,
                        'erro_request': False,
                        'df': pd.DataFrame()
                    }
            else:
                return {
                    'cod_empresa_principal': empresa_principal.cod,
                    'nome_empresa_principal': empresa_principal.nome,
                    'id_empresa': empresa.id_empresa,
                    'nome_empresa': empresa.razao_social,
                    'status': f"erro soc: {response['msg_erro']}",
                    'erro_soc': True,
                    'erro_request': False,
                    'df': pd.DataFrame()
                }
        else:
            return {
                'cod_empresa_principal': empresa_principal.cod,
                'nome_empresa_principal': empresa_principal.nome,
                'id_empresa': empresa.id_empresa,
                'nome_empresa': empresa.razao_social,
                'status': 'erro no request',
                'erro_soc': False,
                'erro_request': True,
                'df': pd.DataFrame()
            }

    # RELATORIOS --------------------------------------------------------------------
    @classmethod
    def rotina_absenteismo(
        self,
        cod_empresa_principal: int,
        id_empresa: int,
        nome_empresa: str,
        data_inicio: datetime,
        data_fim: datetime,
        emails_destinatario: list[str],
        corpo_email: str,
        id_unidade: int = None,
        nome_unidade: str = None,
        testando: bool = False
    ) -> tuple:
        """
        Gera os relatorios Excel e PPT, zipa e envia por email.

        Gera ppt apenas se o DataFrame tiver mais de 50 linhas

        Returns:
            tuple: id_empresa, nome_empresa, id_unidade, nome_unidade, \
            emails_destinatario, len(df), status
        """
        query = self.buscar_licencas(
            cod_empresa_principal=cod_empresa_principal,
            id_empresa=id_empresa,
            id_unidade=id_unidade,
            data_inicio=data_inicio,
            data_fim=data_fim
        )

        df = pd.read_sql(sql=query.statement, con=database.session.bind)

        if nome_unidade:
            caminho_arqvs = f"{UPLOAD_FOLDER}/Absenteismo_{secure_filename(nome_unidade).replace('.', '_')}_{int(datetime.now().timestamp())}"
        else:
            caminho_arqvs = f"{UPLOAD_FOLDER}/Absenteismo_{secure_filename(nome_empresa).replace('.', '_')}_{int(datetime.now().timestamp())}"

        if not df.empty:
            df['hora_inicio_licenca'] = pd.to_datetime(df['hora_inicio_licenca'], unit='s').dt.strftime('%H:%M')
            df['hora_fim_licenca'] = pd.to_datetime(df['hora_fim_licenca'], unit='s').dt.strftime('%H:%M')

            df2 = df[self.COLUNAS_PLANILHA]
            nome_excel = f'{caminho_arqvs}.xlsx'
            df2.to_excel(nome_excel, index=False, freeze_panes=(1, 0))
            arquivos_zipar: list[str] = [nome_excel]

            if id_unidade:
                qtd_ativos = (
                    database.session.query(Funcionario)
                    .filter(Funcionario.situacao == 'Ativo')
                    .filter(Funcionario.id_unidade == id_unidade)
                    .count()
                )
            else:
                qtd_ativos = (
                    database.session.query(Funcionario)
                    .filter(Funcionario.situacao == 'Ativo')
                    .filter(Funcionario.id_empresa == id_empresa)
                    .count()
                )

            if len(df) >= 50:
                nome_ppt = f'{caminho_arqvs}.pptx'
                arquivos_zipar.append(nome_ppt)
                self.criar_ppt(
                    df=df,
                    funcionarios_ativos=qtd_ativos,
                    nome_arquivo=nome_ppt,
                    nome_empresa=nome_empresa,
                    nome_unidade=nome_unidade
                )

            pasta_zip = zipar_arquivos(
                caminhos_arquivos=arquivos_zipar,
                caminho_pasta_zip=f'{caminho_arqvs}.zip'
            )

            if nome_unidade:
                ass = f"Absenteísmo Unidade - {nome_unidade}"
            else:
                ass = f"Absenteísmo Empresa - {nome_empresa}"

            if testando:
                enviar_para = ['gabrielsantos@grsnucleo.com.br']
            else:
                enviar_para = emails_destinatario

            try:
                EmailConnect.send_email(
                    to_addr=enviar_para,
                    message_subject=ass,
                    message_body=corpo_email,
                    message_imgs=[EmailConnect.ASSINATURA_BOT],
                    message_attachments=[pasta_zip],
                    reply_to=['gabrielsantos@grsnucleo.com.br']
                )
                
                # registrar envio
                log_email = EmailConnect(
                    email_to = ','.join(enviar_para),
                    email_subject = ass,
                    attachments = pasta_zip,
                    status = True,
                    df_len = len(df),
                    email_date = datetime.now(tz=TIMEZONE_SAO_PAULO)
                )
                database.session.add(log_email)
                database.session.commit()
                
                return (
                    id_empresa,
                    nome_empresa,
                    id_unidade,
                    nome_unidade,
                    emails_destinatario,
                    len(df),
                    'OK'
                )
            except Exception as erro:
                log_email = EmailConnect(
                    email_to = ','.join(enviar_para),
                    email_subject = ass,
                    attachments = pasta_zip,
                    status = False,
                    error = type(erro).__name__,
                    df_len = len(df),
                    email_date = datetime.now(tz=TIMEZONE_SAO_PAULO)
                )
                database.session.add(log_email)
                database.session.commit()

                return (
                    id_empresa,
                    nome_empresa,
                    id_unidade,
                    nome_unidade,
                    emails_destinatario,
                    len(df),
                    'ERRO'
                )
        else:
            return (
                id_empresa,
                nome_empresa,
                id_unidade,
                nome_unidade,
                emails_destinatario,
                0,
                'NULL'
            )

    @classmethod
    def criar_ppt(
        self,
        df: pd.DataFrame,
        funcionarios_ativos: int,
        nome_arquivo: str,
        nome_empresa: str,
        nome_unidade: str = None
    ):
        df = df.copy()
        # instanciar apresentacao
        presentation = Presentation(self.BASE_PPT)


        # SLIDE 1  - titulo--------------------------------------------------
        slide = presentation.slides[0]

        # periodo
        inicio = df['data_ficha'].min().strftime('%d/%m/%Y')
        fim = df['data_ficha'].max().strftime('%d/%m/%Y')
        slide.shapes[4].text_frame.paragraphs[0].runs[-1].text = f'{inicio} - {fim}'

        # nome empresa
        slide.shapes[5].text_frame.paragraphs[0].runs[-1].text = f'{nome_empresa}'

        # nome unidade
        if nome_unidade:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = f'{nome_unidade}'
        else:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = 'Todas'


        # SLIDE 2 - cards --------------------------------------------------------
        slide = presentation.slides[1]
        shapes = slide.shapes

        # total funcionarios ativos
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = str(funcionarios_ativos)

        # total de licencas
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df))

        # total dias afastados
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(df['dias_afastado'].sum())

        # total func afastados
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df['id_funcionario'].drop_duplicates()))

        # licencas > 15 dias
        shapes[7].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df[df['dias_afastado'] > 15]))

        # taxa absenteismo: horas perdidas / horas previstas
        # shapes[14].text_frame.paragraphs[0].runs[0].text = str()

        # media licencas/funcionario (IFL qtd de licencas / qtd funcionarios ativos)
        shapes[8].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(len(df) / len(df['id_funcionario'].drop_duplicates()))

        # media dias afast/funcionario (IG dias afastados / qtd funcionarios ativos)
        shapes[9].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(df['dias_afastado'].sum() / len(df['id_funcionario'].drop_duplicates()))

        # media dias por licenca (IDA dias afastados / qtd de funcionarios afastados)
        shapes[10].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(df['dias_afastado'].mean(skipna=True))


        # SLIDE 3 ----------------------------------------------------
        slide = presentation.slides[2]
        shapes = slide.shapes

        df['data_ficha'] = pd.to_datetime(df['data_ficha'])
        df['mes'] = df['data_ficha'].dt.month
        df['ano'] = df['data_ficha'].dt.year
        df['mesAno'] = df['data_ficha'].dt.strftime('%m/%Y')

        qtds = (
            df[['id_funcionario']]
            .groupby(by=[
                df['ano'],
                df['mes'],
                df['mesAno']
            ])
            .count()
        )

        qtds.reset_index(inplace=True, level=['mes', 'ano'], drop=True)

        # media licencas/mes
        media = qtds.mean(skipna=True).values[0]
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)

        # min licencas/mes
        minimo = qtds.min(skipna=True).values[0]
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # max licencas/mes
        maximo = qtds.max(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        self.editar_grafico(
            chart=slide.shapes[6].chart,
            categorias=list(qtds.index),
            series={'qtd': list(qtds['id_funcionario'].values)}
        )


        # SLIDE 4 ----------------------------------------------------
        slide = presentation.slides[3]
        shapes = slide.shapes

        qtds = (
            df[['id_funcionario']]
            .drop_duplicates(subset='id_funcionario')
            .groupby(by=[
                df['ano'],
                df['mes'],
                df['mesAno']
            ])
            .count()
        )

        qtds.reset_index(inplace=True, level=['mes', 'ano'], drop=True)

        # media  funcionarios/mes
        media = qtds.mean(skipna=True).values[0]
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)

        # min funcionarios/mes
        minimo = qtds.min(skipna=True).values[0]
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # max funcionarios/mes
        maximo = qtds.max(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        self.editar_grafico(
            chart=slide.shapes[6].chart,
            categorias=list(qtds.index),
            series={'qtd': list(qtds['id_funcionario'].values)}
        )


        # SLIDE 5 ----------------------------------------------------
        slide = presentation.slides[4]
        shapes = slide.shapes

        # qtd licencas por dia da semana
        dias = {
            0: 'Seg',
            1: 'Ter',
            2: 'Qua',
            3: 'Qui',
            4: 'Sex',
            5: 'Sáb',
            6: 'Dom'
        }
        df['dia_da_sem'] = df['data_ficha'].dt.dayofweek
        df = df.replace(to_replace={'dia_da_sem': dias})

        dados = df['dia_da_sem'].value_counts()
        indice = {
            'Dom': 0,
            'Seg': 1,
            'Ter': 2,
            'Qua': 3,
            'Qui': 4,
            'Sex': 5,
            'Sáb': 6
        }
        dados = dados.sort_index(key=lambda x: x.map(indice))

        self.editar_grafico(
            chart=shapes[3].chart,
            categorias=list(dados.index),
            series={'qtd': list(dados.values)}
        )

        # qtd licencas por cid
        # remover cids vazios para evitar contagem de valores NaN
        dados = df.dropna(subset=['cid'])
        dados['grupo_cid'] = dados['cid'].astype(str).str[0].str.upper()
        dados = dados['grupo_cid'].value_counts().sort_values(ascending=False)[:10]
        self.editar_grafico(
            chart=shapes[4].chart,
            categorias=list(dados.index),
            series={'qtd': list(dados.values)}
        )

        # qtd licencas por dias
        dados = df['dias_afastado'].value_counts().sort_values(ascending=False)[:10]
        self.editar_grafico(
            chart=shapes[5].chart,
            categorias=list(dados.index),
            series={'qtd': list(dados.values)}
        )

        # licencas abonadas
        dados = df['abonado'].value_counts().reset_index()
        dados = dados.replace({True: 'Sim', False: 'Não'})
        dados['abonado'] = dados['abonado'] / len(df) # calc percentuais
        dados.sort_values(by='abonado', ascending=False, inplace=True)
        self.editar_grafico(
            chart=shapes[6].chart,
            categorias=list(dados['index'].values),
            series={'Abonado': list(dados['abonado'].values)}
        )


        # # SLIDE 6-8 RANKINGS  ----------------------------------------
        infos_slides_tabelas = [
            (5, 'nome_unidade'),
            (6, 'nome_setor'),
            (7, 'nome_cargo'),
            (8, 'nome_funcionario')
        ]
        
        # adicionar col para usar na soma de licencas
        df['qtd_licencas'] = 1

        for num_slide, nome_col in infos_slides_tabelas:
            slide = presentation.slides[num_slide]
            shapes = slide.shapes
            
            dados = (
                df[[nome_col, 'qtd_licencas', 'dias_afastado']]
                .groupby(nome_col)
                .sum()
                .fillna(0)
            )
            dados = dados.reset_index()
            dados = dados.sort_values('dias_afastado', ascending=False)

            self.editar_tabela(
                shape=shapes[3].table,
                df=dados[:10]
            )


        # SLIDE 10 referecias ------------------------------------------------------
        slide = presentation.slides[9]
        shapes = slide.shapes

        # relatorio gerado em
        slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = datetime.now().strftime('%d/%m/%Y')

        presentation.save(nome_arquivo)
        return None

    @staticmethod
    def editar_grafico(
        chart,
        categorias: list,
        series: dict[list]
    ):
        chart_data = ChartData()
        if series and categorias:
            chart_data.categories = categorias
            for nome_serie, vals in series.items():
                chart_data.add_series(nome_serie, vals)
        else:
            chart_data.categories = ['Sem dados']
            chart_data.add_series('Sem dados', [0])
        chart.replace_data(chart_data)
        return None

    @staticmethod
    def editar_tabela(
        shape: object,
        df: pd.DataFrame,
        font_size: int=14
    ):
        for coluna in range(len(df.columns)):
            # pular duas primeiras celulas de cada coluna
            dados_coluna = [None, None] + list(df[df.columns[coluna]].values)
            for linha in range(2, len(dados_coluna)):
                shape.cell(linha, coluna).text_frame.paragraphs[0].text = str(dados_coluna[linha])
                shape.cell(linha, coluna).text_frame.paragraphs[0].runs[0].font.size = Pt(font_size)
        return None

