from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.util import Pt
from werkzeug.utils import secure_filename

from manager import TIMEZONE_SAO_PAULO, UPLOAD_FOLDER, database
from manager.email_connect import EmailConnect
from manager.models import (Empresa, EmpresaPrincipal, Exame, Prestador,
                            TipoExame, Unidade, Funcionario)
from manager.utils import get_json_configs, zipar_arquivos


class ExamesRealizados(database.Model):
    __tablename__ = 'ExamesRealizados'
    id_realizado = database.Column(database.Integer, primary_key=True)
    cod_empresa_principal = database.Column(database.Integer, database.ForeignKey('EmpresaPrincipal.cod'), nullable=False)
    seq_ficha = database.Column(database.Integer, nullable=False)
    cod_funcionario = database.Column(database.Integer, nullable=False)
    cpf = database.Column(database.String(30))
    nome_funcionario = database.Column(database.String(150))
    data_criacao = database.Column(database.Date, nullable=False)
    data_ficha = database.Column(database.Date, nullable=False)
    data_resultado = database.Column(database.Date, nullable=False)
    
    cod_tipo_exame = database.Column(database.Integer, database.ForeignKey('TipoExame.cod_tipo_exame'), nullable=False)
    id_prestador = database.Column(database.Integer, database.ForeignKey('Prestador.id_prestador'))
    id_empresa = database.Column(database.Integer, database.ForeignKey('Empresa.id_empresa'), nullable=False)
    id_unidade = database.Column(database.Integer, database.ForeignKey('Unidade.id_unidade'), nullable=False)
    id_exame = database.Column(database.Integer, database.ForeignKey('Exame.id_exame'), nullable=False)

    cod_setor = database.Column(database.String(255))
    nome_setor = database.Column(database.String(255))
    cod_cargo = database.Column(database.String(255))
    nome_cargo = database.Column(database.String(255))

    BASE_PPT = 'manager/modules/exames_realizados/ppt/base_exames_realizados_v2.pptx'

    colunas_planilha = [
        'cod_empresa_principal',
        'nome',
        'cod_empresa',
        'razao_social',
        'cod_unidade',
        'nome_unidade',
        'cod_setor',
        'nome_setor',
        'cod_cargo',
        'nome_cargo',
        'cod_funcionario',
        'cpf',
        'nome_funcionario',
        'seq_ficha',
        'cod_exame',
        'nome_exame',
        'cod_tipo_exame',
        'nome_tipo_exame',
        'data_criacao',
        'data_ficha',
        'data_resultado',
        'cod_prestador',
        'nome_prestador'
    ]

    @classmethod
    def buscar_exames_realizados(
        self,
        cod_empresa_principal: int,
        id_empresa: int = None,
        id_unidade: int = None,
        id_exame: int = None,
        data_inicio: datetime = None,
        data_fim: datetime = None
    ):

        filtros = [(self.cod_empresa_principal == cod_empresa_principal)]

        if id_empresa:
            filtros.append(self.id_empresa == id_empresa)
        if id_exame:
            filtros.append(self.id_exame == id_exame)
        if data_inicio:
            filtros.append(self.data_resultado >= data_inicio)
        if data_fim:
            filtros.append(self.data_resultado <= data_fim)
        if id_unidade:
            filtros.append(self.id_unidade == id_unidade)
        
        query = (
            database.session.query(
                self,
                EmpresaPrincipal.nome,
                Empresa.cod_empresa,
                Empresa.razao_social,
                Unidade.cod_unidade,
                Unidade.nome_unidade,
                Prestador.cod_prestador,
                Prestador.nome_prestador,
                TipoExame.cod_tipo_exame,
                TipoExame.nome_tipo_exame,
                Exame.cod_exame,
                Exame.nome_exame
            )
            .outerjoin(EmpresaPrincipal, self.cod_empresa_principal == EmpresaPrincipal.cod)
            .outerjoin(Empresa, self.id_empresa == Empresa.id_empresa)
            .outerjoin(Unidade, self.id_unidade == Unidade.id_unidade)
            .outerjoin(Prestador, self.id_prestador == Prestador.id_prestador)
            .outerjoin(TipoExame, self.cod_tipo_exame == TipoExame.cod_tipo_exame)
            .outerjoin(Exame, self.id_exame == Exame.id_exame)
            .filter(*filtros)
        )

        return query

    @classmethod
    def inserir_exames_realizados(
        self,
        id_empresa: int,
        dataInicio: datetime,
        dataFim: datetime
    ) -> int:
        """
        Realiza request para o Exporta Dados Pedido Exame \
        e insere novos exames realizados na Tabela. Remove \
        exames ja realizados baseando no seq_ficha + id_exame.

        Returns:
            dict: {
                'cod_empresa_principal': int,
                'nome_empresa_principal': int,
                'id_empresa': int,
                'nome_empresa': str,
                'status': str,
                'qtd': int
            }
        """
        from manager.exporta_dados import ExportaDadosWS

        empresa: Empresa = Empresa.query.get(id_empresa)
        empresa_principal: EmpresaPrincipal = EmpresaPrincipal.query.get(empresa.cod_empresa_principal)
        credenciais: dict = get_json_configs(empresa_principal.configs_exporta_dados)

        par: dict = ExportaDadosWS.pedido_exame(
            empresa = empresa.cod_empresa,
            cod_exporta_dados = credenciais['EXPORTADADOS_PEDIDOEXAME_COD'],
            chave = credenciais['EXPORTADADOS_PEDIDOEXAME_KEY'],
            dataInicio = dataInicio.strftime('%d/%m/%Y'),
            dataFim = dataFim.strftime('%d/%m/%Y')
        )

        response: dict = ExportaDadosWS.request_exporta_dados_ws(
            parametro=par,
            id_empresa=empresa.id_empresa,
            obs='Pedido Exame'
        )

        if response['response'].status_code == 200:
            if not response['erro_soc']:
                response_text: str = response['response'].text
                df = ExportaDadosWS.xml_to_dataframe(xml_string=response_text)

                if not df.empty:
                    df = df.replace({'': None})
                    # manter apenas exames com data de resultado
                    df = df[~df['DATAEXAME'].isna()]
                    if not df.empty:
                        # pegar id do Exame
                        df_database = pd.read_sql(
                            sql=(
                                database.session.query(
                                    Exame.id_exame,
                                    Exame.cod_exame,
                                )
                                .filter(Exame.cod_empresa_principal == empresa_principal.cod)
                                .statement
                            ),
                            con=database.session.bind
                        )
                        df = pd.merge(
                            df,
                            df_database,
                            how='left',
                            left_on='CODIGOINTERNOEXAME',
                            right_on='cod_exame'
                        ) 

                        df['id_empresa'] = empresa.id_empresa

                        # pegar id unidade
                        df_database = pd.read_sql(
                            sql=(
                                database.session.query(
                                    Unidade.id_unidade,
                                    Unidade.cod_unidade
                                )
                                .filter(Unidade.id_empresa == empresa.id_empresa)
                                .statement
                            ),
                            con=database.session.bind
                        )
                        df = pd.merge(
                            df,
                            df_database,
                            how='left',
                            left_on='CODIGOUNIDADE',
                            right_on='cod_unidade'
                        ) 

                        # pegar id do Prestador
                        df_database = pd.read_sql(
                            sql=(
                                database.session.query(
                                    Prestador.id_prestador,
                                    Prestador.cod_prestador
                                )
                                .filter(Prestador.cod_empresa_principal == empresa_principal.cod)
                                .statement
                            ),
                            con=database.session.bind
                        )
                        # adicionar uma linha Nan no df database para nao travar o merge quando o prestador \
                        # estiver vazio no df
                        df_database = pd.concat(
                            [df_database, pd.DataFrame({'id_prestador': [None], 'cod_prestador': [None]})],
                            axis=0,
                            ignore_index=True
                        )
                        df_database = df_database.astype('Int32')
                        df['CODIGOPRESTADOR'] = df['CODIGOPRESTADOR'].astype('Int32')
                        df = pd.merge(
                            df,
                            df_database,
                            how='left',
                            left_on='CODIGOPRESTADOR',
                            right_on='cod_prestador'
                        ) 

                        # tratar colunas
                        for col in ['SEQUENCIAFICHA', 'CODIGOFUNCIONARIO', 'CODIGOTIPOEXAME', 'CODIGOEMPRESA']:
                            df[col] = df[col].astype(int)

                        for col in ['DATACRIACAOPEDIDOEXAMES', 'DATAFICHA', 'DATAEXAME']:
                            df[col] = pd.to_datetime(df[col], dayfirst=True, errors='coerce').dt.date
                        # dropar linhas com erro nas datas
                        df.dropna(axis=0, subset=['DATACRIACAOPEDIDOEXAMES', 'DATAFICHA', 'DATAEXAME'], inplace=True)

                        df['CPFFUNCIONARIO'] = df['CPFFUNCIONARIO'].astype('string')

                        df['cod_empresa_principal'] = empresa_principal.cod

                        df = df.rename(columns={
                            'SEQUENCIAFICHA': 'seq_ficha',
                            'CODIGOFUNCIONARIO': 'cod_funcionario',
                            'CPFFUNCIONARIO': 'cpf',
                            'NOMEFUNCIONARIO': 'nome_funcionario',
                            'DATACRIACAOPEDIDOEXAMES': 'data_criacao',
                            'DATAFICHA': 'data_ficha',
                            'DATAEXAME': 'data_resultado',
                            'CODIGOTIPOEXAME': 'cod_tipo_exame',
                            'CODIGOSETOR': 'cod_setor',
                            'NOMESETOR': 'nome_setor',
                            'CODIGOCARGO': 'cod_cargo',
                            'NOMECARGO': 'nome_cargo'
                        })

                        df.dropna(
                            axis=0,
                            subset=['id_empresa', 'id_unidade'],
                            inplace=True
                        )

                        df = df[[
                            'cod_empresa_principal',
                            'seq_ficha',
                            'cod_funcionario',
                            'cpf',
                            'nome_funcionario',
                            'data_criacao',
                            'data_ficha',
                            'data_resultado',
                            'cod_tipo_exame',
                            'id_prestador',
                            'id_empresa',
                            'id_unidade',
                            'id_exame',
                            'cod_setor',
                            'nome_setor',
                            'cod_cargo',
                            'nome_cargo'
                        ]]

                        # remover realizados que ja existem
                        df_database = pd.read_sql(
                            sql=(
                                database.session.query(
                                    ExamesRealizados.id_realizado,
                                    ExamesRealizados.seq_ficha,
                                    ExamesRealizados.id_exame,
                                )
                                .filter(ExamesRealizados.id_empresa == empresa.id_empresa)
                                .statement
                            ),
                            con=database.session.bind
                        )
                        df = pd.merge(
                            df,
                            df_database,
                            how='left',
                            on=['seq_ficha', 'id_exame']
                        )
                        df = df[df['id_realizado'].isna()]

                        qtd_inseridos = df.to_sql(name=self.__tablename__, con=database.session.bind, if_exists='append', index=False)
                        database.session.commit()
                        return {
                            'cod_empresa_principal': empresa_principal.cod,
                            'nome_empresa_principal': empresa_principal.nome,
                            'id_empresa': empresa.id_empresa,
                            'nome_empresa': empresa.razao_social,
                            'status': 'ok',
                            'qtd': qtd_inseridos
                        }
                    else:
                        return {
                            'cod_empresa_principal': empresa_principal.cod,
                            'nome_empresa_principal': empresa_principal.nome,
                            'id_empresa': empresa.id_empresa,
                            'nome_empresa': empresa.razao_social,
                            'status': 'vazio',
                            'qtd': 0
                        }
                else:
                    return {
                            'cod_empresa_principal': empresa_principal.cod,
                            'nome_empresa_principal': empresa_principal.nome,
                            'id_empresa': empresa.id_empresa,
                            'nome_empresa': empresa.razao_social,
                            'status': 'vazio',
                            'qtd': 0
                        }
            else:
                return {
                            'cod_empresa_principal': empresa_principal.cod,
                            'nome_empresa_principal': empresa_principal.nome,
                            'id_empresa': empresa.id_empresa,
                            'nome_empresa': empresa.razao_social,
                            'status': f"erro soc: {response['msg_erro']}",
                            'qtd': 0
                        }
        else:
            return {
                        'cod_empresa_principal': empresa_principal.cod,
                        'nome_empresa_principal': empresa_principal.nome,
                        'id_empresa': empresa.id_empresa,
                        'nome_empresa': empresa.razao_social,
                        'status': 'erro no request',
                        'qtd': 0
                    }

    @classmethod
    def rotina_exames_realizados(
        self,
        cod_empresa_principal: int,
        id_empresa: int,
        nome_empresa: str,
        data_inicio: datetime,
        data_fim: datetime,
        emails_destinatario: list[str],
        corpo_email: str,
        id_unidade: int = None,
        nome_unidade: str = None,
        filtro_exames: list[str] = None,
        testando: bool = False
    ) -> tuple:
        """
        Args:
            filtro_exames (list[str], optional): lista de ids dos Exames \
            para filtrar na query. Defaults to None.

            Gera ppt apenas se a qtd de linhas for maior que 50

        Returns:
            tuple[int, str]: (num de linhas da tabela, status email)
        """
        query = ExamesRealizados.buscar_exames_realizados(
            cod_empresa_principal=cod_empresa_principal,
            id_empresa=id_empresa,
            id_unidade=id_unidade,
            data_inicio=data_inicio,
            data_fim=data_fim
        )

        if filtro_exames:
            query = query.filter(Exame.id_exame.in_(filtro_exames))
        
        df = pd.read_sql(sql=query.statement, con=database.session.bind)

        if nome_unidade:
            caminho_arqvs = f"{UPLOAD_FOLDER}/ExamesRealizados_{secure_filename(nome_unidade).replace('.', '_')}_{int(datetime.now().timestamp())}"
        else:
            caminho_arqvs = f"{UPLOAD_FOLDER}/ExamesRealizados_{secure_filename(nome_empresa).replace('.', '_')}_{int(datetime.now().timestamp())}"

        if not df.empty:
            df2 = df[ExamesRealizados.colunas_planilha]
            nome_excel = f'{caminho_arqvs}.xlsx'
            df2.to_excel(nome_excel, index=False, freeze_panes=(1, 0))
            arquivos_zipar: list[str] = [nome_excel]

            if len(df) > 50:
                nome_ppt = f'{caminho_arqvs}.pptx'
                arquivos_zipar.append(nome_ppt)
                ExamesRealizados.criar_ppt(
                    df=df,
                    nome_arquivo=nome_ppt,
                    nome_empresa=nome_empresa,
                    nome_unidade=nome_unidade
                )

            pasta_zip = zipar_arquivos(
                caminhos_arquivos=arquivos_zipar,
                caminho_pasta_zip=f'{caminho_arqvs}.zip'
            )

            if nome_unidade:
                ass = f"Exames Realizados Unidades - {nome_unidade}"
            else:
                ass = f"Exames Realizados Empresas - {nome_empresa}"

            if testando:
                enviar_para = ['gabrielsantos@grsnucleo.com.br']
            else:
                enviar_para = emails_destinatario

            try:
                EmailConnect.send_email(
                    to_addr=enviar_para,
                    message_subject=ass,
                    message_body=corpo_email,
                    message_imgs=[EmailConnect.ASSINATURA_BOT],
                    message_attachments=[pasta_zip],
                    reply_to=['gabrielsantos@grsnucleo.com.br']
                )
                
                # registrar envio
                log_email = EmailConnect(
                    email_to = ','.join(enviar_para),
                    email_subject = ass,
                    attachments = pasta_zip,
                    status = True,
                    df_len = len(df),
                    email_date = datetime.now(tz=TIMEZONE_SAO_PAULO)
                )
                database.session.add(log_email)
                database.session.commit()
                
                return (
                    id_empresa,
                    nome_empresa,
                    id_unidade,
                    nome_unidade,
                    emails_destinatario,
                    len(df),
                    'OK'
                )
            except Exception as erro:
                log_email = EmailConnect(
                    email_to = ','.join(enviar_para),
                    email_subject = ass,
                    attachments = pasta_zip,
                    status = False,
                    error = type(erro).__name__,
                    df_len = len(df),
                    email_date = datetime.now(tz=TIMEZONE_SAO_PAULO)
                )
                database.session.add(log_email)
                database.session.commit()
                
                return (
                    id_empresa,
                    nome_empresa,
                    id_unidade,
                    nome_unidade,
                    emails_destinatario,
                    len(df),
                    'ERRO'
                )
        else:
            return (
                id_empresa,
                nome_empresa,
                id_unidade,
                nome_unidade,
                emails_destinatario,
                0,
                'NULL'
            )

    @classmethod
    def criar_ppt(
        self,
        df: pd.DataFrame,
        nome_arquivo: str,
        nome_empresa: str,
        nome_unidade: str = None,
        data_origem: datetime = datetime.now()
    ):
        # NOTE: ao editar um grafico de barras placeholder com varias series \
        # no Powerpoint, inserir series extras para evitar que o python crie \
        # series novas, assim as cores seguidas pelas series serao as do \
        # tema do PPT. Quando faltam series, as novas sao criadas com a mesma cor \
        # mas se tem series a mais, as que sobraram saem do PPT final

        df = df.copy()

        # instanciar apresentacao
        presentation = Presentation(self.BASE_PPT)


        # SLIDE 1  - titulo--------------------------------------------------
        slide = presentation.slides[0]

        # periodo
        inicio = df['data_resultado'].min().strftime('%d/%m/%Y')
        fim = df['data_resultado'].max().strftime('%d/%m/%Y')
        slide.shapes[4].text_frame.paragraphs[0].runs[-1].text = f'{inicio} - {fim}'

        # nome empresa
        slide.shapes[5].text_frame.paragraphs[0].runs[-1].text = f'{nome_empresa}'

        # nome unidade
        if nome_unidade:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = f'{nome_unidade}'
        else:
            slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = 'Todas'


        # SLIDE 2 - indicadores gerais-------------------------------------------
        slide = presentation.slides[1]
        shapes = slide.shapes

        # total exames
        shapes[3].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df))

        # total pedidos
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = str(len(df['seq_ficha'].drop_duplicates()))

        qtds = df[['seq_ficha', 'nome_exame']].groupby(by='seq_ficha').count()
        # media exames/pedido
        media = qtds.mean(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)

        # max exames/pedido
        maximo = qtds.max(skipna=True).values[0]
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        # min exames/pedido
        minimo = qtds.min(skipna=True).values[0]
        shapes[7].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # media exames/unidade
        media = df[['cod_unidade', 'nome_exame']].groupby(by='cod_unidade').count().mean(skipna=True).values[0]
        shapes[8].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)

        # media exames/setor
        media = df[['cod_setor', 'nome_exame']].groupby(by='cod_setor').count().mean(skipna=True).values[0]
        shapes[9].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)

        # media exames/cargo
        media = df[['cod_cargo', 'nome_exame']].groupby(by='cod_cargo').count().mean(skipna=True).values[0]
        shapes[10].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)

        # media exames/funcionario
        media = df[['cod_funcionario', 'nome_exame']].groupby(by='cod_funcionario').count().mean(skipna=True).values[0]
        shapes[11].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)


        # SLIDE 3 - Exames por data---------------------------------------
        slide = presentation.slides[2]
        shapes = slide.shapes

        df['data_resultado'] = pd.to_datetime(df['data_resultado'])
        df['mesAno'] = df['data_resultado'].dt.strftime('%m/%Y')

        qtds = df[['mesAno', 'nome_exame']].groupby(by='mesAno').count()

        # media exames/data
        media = qtds.mean(skipna=True).values[0]
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)

        # min exames/data
        minimo = qtds.min(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # max exames/data
        maximo = qtds.max(skipna=True).values[0]
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        df['mes'] = df['data_resultado'].dt.month
        df['ano'] = df['data_resultado'].dt.year

        dados = (
            df[['ano', 'mes', 'mesAno', 'nome_tipo_exame', 'nome_exame']]
            .pivot_table(
                values='nome_exame',
                index=['ano', 'mes', 'mesAno'],
                columns='nome_tipo_exame',
                aggfunc='count',
            )
            .fillna(0)
        )
        dados.reset_index(level=['ano', 'mes'], drop=True, inplace=True)
        # organizar colunas para manter a ordem no grafico
        dados.loc['Total'] = dados.sum(numeric_only=True)
        dados = dados.sort_values(by='Total', axis=1, ascending=False)
        dados.drop(labels=['Total'], axis=0, inplace=True)

        self.editar_grafico(
            chart=slide.shapes[3].chart,
            categorias=list(dados.index),
            series={col: list(dados[col].values) for col in dados.columns}
        )


        # SLIDE 4 - exames por tipo ------------------------------------------------------
        slide = presentation.slides[3]
        shapes = slide.shapes

        qtds = df[['nome_tipo_exame', 'nome_exame']].groupby(by='nome_tipo_exame').count()
        # media
        media = qtds.mean(skipna=True).values[0]
        shapes[4].shapes[1].text_frame.paragraphs[0].runs[0].text = '{:.2f}'.format(media)

        # min
        minimo = qtds.min(skipna=True).values[0]
        shapes[5].shapes[1].text_frame.paragraphs[0].runs[0].text = str(minimo)

        # max
        maximo = qtds.max(skipna=True).values[0]
        shapes[6].shapes[1].text_frame.paragraphs[0].runs[0].text = str(maximo)

        dados = df[['nome_tipo_exame', 'nome_exame']].groupby('nome_tipo_exame').count()
        dados['nome_exame'] = dados['nome_exame'] / dados['nome_exame'].sum()
        dados = dados.sort_values('nome_exame', axis=0, ascending=False)

        self.editar_grafico(
            chart=slide.shapes[3].chart,
            categorias=list(dados.index),
            series={'qtd': list(dados['nome_exame'].values)}
        )


        # SLIDES DE TABELA 5-10 ------------------------------------------------------
        infos_slides_tabelas = [
            (4, 'nome_unidade'),
            (5, 'nome_setor'),
            (6, 'nome_cargo'),
            (7, 'nome_exame'),
            (8, 'nome_prestador'),
            (9, 'nome_funcionario')
        ]

        for num_slide, nome_col in infos_slides_tabelas:
            slide = presentation.slides[num_slide]
            shapes = slide.shapes
            
            dados = df[[nome_col, 'razao_social']].groupby(nome_col).count().fillna(0)
            dados = dados.reset_index()
            dados = dados.sort_values('razao_social', ascending=False)

            self.editar_tabela(
                shape=shapes[3].table,
                df=dados[:10]
            )

        # SLIDE 11 referecias ------------------------------------------------------
        slide = presentation.slides[10]
        shapes = slide.shapes

        # relatorio gerado em
        slide.shapes[6].text_frame.paragraphs[0].runs[-1].text = data_origem.strftime('%d/%m/%Y')

        presentation.save(nome_arquivo)
        return None

    @staticmethod
    def editar_grafico(
        chart,
        categorias: list,
        series: dict[list]
    ):
        chart_data = ChartData()
        chart_data.categories = categorias
        if series:
            for nome_serie, vals in series.items():
                chart_data.add_series(nome_serie, vals)
        else:
            chart_data.add_series('Sem dados', [0])
        chart.replace_data(chart_data)
        return None

    @staticmethod
    def editar_tabela(
        shape: object,
        df: pd.DataFrame,
        font_size: int=14
    ):
        for coluna in range(len(df.columns)):
            # pular duas primeiras celulas de cada coluna
            dados_coluna = [None, None] + list(df[df.columns[coluna]].values)
            for linha in range(2, len(dados_coluna)):
                shape.cell(linha, coluna).text_frame.paragraphs[0].text = str(dados_coluna[linha])
                shape.cell(linha, coluna).text_frame.paragraphs[0].runs[0].font.size = Pt(font_size)
        return None

